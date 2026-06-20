"""
Warrior Prep — Teacher PowerPoint Builder (Full Content Version)
Every lesson has: intro, reading slides, comprehension questions, teaching content,
activity slides (content on screen), and exit ticket.
All slides have teacher notes with talk tracks.
"""

import os, textwrap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from content import UNITS, SCHOOL_NAME, COURSE_TITLE, GRADE
from readings import READINGS
from lesson_meta import BELL_RINGERS, LESSON_TIMING
from stories import STORIES
from youtube_links import YOUTUBE

os.makedirs("output", exist_ok=True)

# BELL_RINGERS and LESSON_TIMING imported from lesson_meta.py
_REMOVED = {
    "1.1": {
        "type": "Launch",
        "title": "Brain Warm-Up",
        "prompt": "Write your answer: In the last 24 hours, name ONE moment when you had to stop yourself from doing something you wanted to do in order to do something you needed to do. What made it hard?",
        "share": "Pair-share (1 min) → cold-call 2 students",
        "connection": "Launches Unit 1 by surfacing a real experience of inhibitory control — the exact EF skill we name and define today.",
    },
    "1.2": {
        "type": "Review",
        "title": "EF Skill Quick-Sort",
        "prompt": "Without looking at your notes — name THREE executive functioning skills from last class. Then: rank them from easiest to hardest to use. Be ready to explain your ranking.",
        "share": "Thumbs up when ready → cold-call 1 student per skill",
        "connection": "Retrieves lesson 1.1 vocabulary. Ranking requires students to apply, not just recall — sets up today's prioritization work.",
    },
    "1.3": {
        "type": "Anticipation",
        "title": "Myth or Fact?",
        "prompt": "Agree or Disagree — write your answer AND your reason:\n'I can study while watching a show and still remember what I studied just as well as if I'd studied in silence.'",
        "share": "Show of hands Agree/Disagree → brief debate (don't confirm yet — the reading will)",
        "connection": "Creates productive cognitive conflict before the lesson on multitasking and working memory.",
    },
    "1.4": {
        "type": "Review",
        "title": "Unit 1 Vocabulary Lightning Round",
        "prompt": "Without your workbook: write the definition of EACH word in your own words.\n1. Executive functioning  2. Inhibitory control  3. Working memory  4. Cognitive flexibility",
        "share": "Partner check — compare definitions and agree on one shared version",
        "connection": "Final vocabulary consolidation before the SMART goal lesson — students will use all four terms in their goal-setting work.",
    },
    "2.1": {
        "type": "Launch",
        "title": "Career Assumption Check",
        "prompt": "Respond quickly — don't overthink:\n'Name one career you think you might be interested in someday. NOW: where does that career fit in school? What courses at Two Rivers connect to it?'",
        "share": "Table share (1 min) → 3 volunteers share out",
        "connection": "Surfaces prior knowledge and reveals what students already understand (or misunderstand) about the connection between school and careers — frames Unit 2.",
    },
    "2.2": {
        "type": "Review",
        "title": "Career Fields Memory Map",
        "prompt": "In 2 minutes: from memory, list as many of the six Minnesota career fields as you can. Under each one, name at least ONE Two Rivers course or program that fits.",
        "share": "Compare with a partner — fill in anything you missed",
        "connection": "Retrieves lesson 2.1 content. Sets up today's focus on engineering/health pathways and the CAPS programs.",
    },
    "2.3": {
        "type": "Anticipation",
        "title": "What Does a Communicator Do?",
        "prompt": "Quick write: What job do you picture when you hear 'Arts, Communications, and Information Systems'? Do you think that field pays well? Do you think it's in demand? Explain your reasoning.",
        "share": "Pair-share → 2 volunteers",
        "connection": "Challenges assumptions about arts/communication careers before the lesson expands what students think is possible in those fields.",
    },
    "2.4": {
        "type": "Review",
        "title": "Two Rivers Pathway Recap",
        "prompt": "Draw a quick web/mind map from memory: put 'Two Rivers HS' in the center. Branch out to every pathway, program, or field we've studied so far. Add at least one course or career under each.",
        "share": "Compare maps with a partner — what's on theirs that's not on yours?",
        "connection": "Synthesizes all of Unit 2 before the final lesson. Gaps in the map show where students need the reading today.",
    },
    "3.1": {
        "type": "Launch",
        "title": "Reading Habits Honest Check",
        "prompt": "Be honest — no one is grading this:\nWhen you read something for school that you don't find interesting, what do you actually do? Check all that apply.\n☐ Read carefully  ☐ Skim it  ☐ Read the first and last line of each paragraph  ☐ Ask someone what it says  ☐ Give up",
        "share": "Quick poll — hands for each option (normalize honesty, no shame)",
        "connection": "Surfaces reading habits without judgment — creates buy-in for today's lesson on active reading strategies.",
    },
    "3.2": {
        "type": "Review",
        "title": "Annotation Symbol Quick-Apply",
        "prompt": "Look at this sentence. Annotate it RIGHT NOW using the symbols from last class:\n'The students who scored highest on the final exam were the ones who had annotated their textbooks most heavily — not the ones who had read the most pages.'",
        "share": "Show your annotations to a partner — did you use the same symbols? Did you mark the same parts?",
        "connection": "Applies lesson 3.1 annotation skills immediately. The sentence content also previews today's lesson on writing stamina.",
    },
    "3.3": {
        "type": "Anticipation",
        "title": "How Do YOU Study?",
        "prompt": "Before you look at any notes — describe your exact study method for a big test:\nStep 1: ___  Step 2: ___  Step 3: ___\nDo you think it works? How would you know?",
        "share": "Partner share — compare methods. Are they similar?",
        "connection": "Activates prior knowledge about study habits before today's lesson challenges students to evaluate those habits against research.",
    },
    "3.4": {
        "type": "Review",
        "title": "Study Strategy Ranking",
        "prompt": "From last class: put these strategies in order from MOST effective to LEAST effective (according to the research, not your preference):\nRe-reading · Highlighting · Retrieval practice · Spaced repetition · Interleaving",
        "share": "Compare rankings with a partner → class consensus",
        "connection": "Retrieves and ranks lesson 3.3 strategies. Sets up today's focus on building independent study advocacy — students need to know the strategies to advocate for time to use them.",
    },
    "4.1": {
        "type": "Launch",
        "title": "What Do You Actually Know About Graduation?",
        "prompt": "Without looking anything up — answer as specifically as you can:\n1. How many total credits do you need to graduate from Two Rivers?\n2. Name one course that is REQUIRED for every student.\n3. What is the difference between required courses and electives?",
        "share": "Table share → reveal answers on next slide (after students commit to answers)",
        "connection": "Reveals what students know (and don't) about graduation before we map it out — creates urgency for careful planning.",
    },
    "4.2": {
        "type": "Review",
        "title": "Graduation Requirement Recall",
        "prompt": "From last class — fill in what you remember:\n• Required English credits in grades 9-12: ___\n• Required math: ___\n• Required science: ___\n• Open elective credits available: ___\n• Warrior Seminar is worth ___ credits and is taken in grade ___",
        "share": "Check with your Academic Planning Guide — how close were you?",
        "connection": "Retrieves the graduation structure before students begin building their actual mock schedule today.",
    },
    "4.3": {
        "type": "Anticipation",
        "title": "Activity Assumption Check",
        "prompt": "Agree, Disagree, or It Depends — write a reason for each:\n1. 'Students who do extracurriculars have lower GPAs because they have less study time.'\n2. 'It's better to try one activity deeply than three activities lightly.'\n3. 'You should decide your activities in 8th grade before 9th starts.'",
        "share": "Pair debate (1 min) — who changed their mind after hearing a partner?",
        "connection": "Surfaces assumptions about involvement before the reading and lesson challenge or confirm them.",
    },
    "4.4": {
        "type": "Review",
        "title": "Rapid-Fire: My 9th Grade Plan So Far",
        "prompt": "Quick-write in 2 minutes — answer all three:\n1. Name TWO courses you've decided to take in 9th grade (from your schedule builder).\n2. Name ONE activity you circled in the activities checklist.\n3. One word that describes how you feel about starting 9th grade.",
        "share": "Popcorn share — one word feelings around the room",
        "connection": "Reconnects to students' personal plans from Lessons 4.1-4.3. The 'one word' share surfaces anxiety early — today's time-map lesson directly addresses it.",
    },
    "5.1": {
        "type": "Launch",
        "title": "Grade vs. GPA — What's the Difference?",
        "prompt": "Without any help — explain the difference between a grade and a GPA.\nThen: if a student gets a B in every class for one semester, what is their GPA? Make your best guess.",
        "share": "Show of hands for different GPA guesses — don't confirm yet",
        "connection": "Surfaces misconceptions about GPA before the lesson and calculator activity — students almost always underestimate how grades translate.",
    },
    "5.2": {
        "type": "Review",
        "title": "GPA Impact Scenario",
        "prompt": "Look at this scenario: A student gets straight B's (3.0) in semester 1. In semester 2, they get one A and the rest B's. By how much does their GPA change?\nThen: what ONE habit change would have the biggest effect on a student's GPA?",
        "share": "Calculate with a partner → share the habit",
        "connection": "Applies lesson 5.1 GPA math before today's lesson on the social-academic skills (BARR, belonging) that protect GPA.",
    },
    "5.3": {
        "type": "Anticipation",
        "title": "If Registration Were TODAY",
        "prompt": "Right now — without looking at anything — write down the 7 courses you would put on your registration form if it were due today.\nPeriod 1: ___  P2: ___  P3: ___  P4: ___  P5: ___  P6: ___  P7: ___",
        "share": "Compare with a partner — what's different? What did you forget?",
        "connection": "Creates urgency and reveals gaps before today's final schedule review and rationale lesson.",
    },
    "5.4": {
        "type": "Review",
        "title": "What I Still Want to Know",
        "prompt": "Think about everything we've covered in Warrior Prep this semester.\nWrite: What is the ONE question about high school that you still haven't gotten a clear answer to? Make it specific. This is your panel question — make it count.",
        "share": "Each student reads their question aloud — class votes (thumbs) if it's a question they ALSO want answered",
        "connection": "Directly generates today's panel question. The vote helps students see which questions are widely shared — builds investment in the panel discussion.",
    },
}

# ── 50-Minute Lesson Timing Plans ────────────────────────────────────────────
# Breakdown for each lesson label showing how 50 minutes is allocated.
LESSON_TIMING = {
    "1.1": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 10, "activity": 10, "exit": 5},
    "1.2": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 8,  "activity": 12, "exit": 5},
    "1.3": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 10, "activity": 10, "exit": 5},
    "1.4": {"bell": 5, "title": 3, "read": 10, "qs": 6, "content": 8,  "activity": 13, "exit": 5},
    "2.1": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 8,  "activity": 12, "exit": 5},
    "2.2": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 10, "activity": 10, "exit": 5},
    "2.3": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 10, "activity": 10, "exit": 5},
    "2.4": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 8,  "activity": 12, "exit": 5},
    "3.1": {"bell": 5, "title": 3, "read": 11, "qs": 7, "content": 9,  "activity": 10, "exit": 5},
    "3.2": {"bell": 5, "title": 3, "read": 10, "qs": 6, "content": 8,  "activity": 13, "exit": 5},
    "3.3": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 10, "activity": 10, "exit": 5},
    "3.4": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 9,  "activity": 11, "exit": 5},
    "4.1": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 8,  "activity": 12, "exit": 5},
    "4.2": {"bell": 5, "title": 3, "read": 9,  "qs": 6, "content": 7,  "activity": 15, "exit": 5},
    "4.3": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 8,  "activity": 12, "exit": 5},
    "4.4": {"bell": 5, "title": 3, "read": 9,  "qs": 6, "content": 8,  "activity": 14, "exit": 5},
    "5.1": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 8,  "activity": 12, "exit": 5},
    "5.2": {"bell": 5, "title": 3, "read": 10, "qs": 7, "content": 10, "activity": 10, "exit": 5},
    "5.3": {"bell": 5, "title": 3, "read": 9,  "qs": 6, "content": 7,  "activity": 15, "exit": 5},
    "5.4": {"bell": 5, "title": 3, "read": 8,  "qs": 5, "content": 4,  "activity": 25, "exit": 5},
}  # end _REMOVED (actual dicts live in lesson_meta.py)

# ── Colors ────────────────────────────────────────────────────────────────────
RED   = RGBColor(0xC7, 0x14, 0x14)
DARK  = RGBColor(0x21, 0x21, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD  = RGBColor(0xF5, 0xE0, 0xB2)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
BLUE  = RGBColor(0xE8, 0xF0, 0xFE)
CREAM = RGBColor(0xFD, 0xFA, 0xF4)
GRAY  = RGBColor(0x88, 0x88, 0x88)


def uc(unit):
    t = unit["color"]
    return RGBColor(int(t[0]*255), int(t[1]*255), int(t[2]*255))


def new_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    s = prs.slides.add_slide(layout)
    for shape in list(s.shapes):
        shape._element.getparent().remove(shape._element)
    return s


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh


def tb(slide, text, x, y, w, h, size=14, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    color = color or DARK
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def mtb(slide, lines, x, y, w, h, size=13, color=None, bold_first=False,
        bullet=False, indent=0.1, spacing=3):
    """Multi-line textbox — each item in lines is a new paragraph."""
    color = color or DARK
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(spacing)
        if bullet:
            p.level = 1
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold_first and (line == lines[0])
        run.font.color.rgb = color
    return box


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header_bar(slide, text, color, y=0, h=0.58):
    rect(slide, 0, y, 10, h, color)
    tb(slide, text, 0.18, y + 0.07, 9.6, h - 0.1, size=17, bold=True, color=WHITE)


def sub_bar(slide, text, color, y, h=0.38):
    rect(slide, 0.2, y, 9.6, h, color)
    tb(slide, text, 0.32, y + 0.04, 9.3, h - 0.08, size=12, bold=True, color=WHITE)


def timer_chip(slide, minutes, x=8.5, y=0.65):
    rect(slide, x, y, 1.35, 0.4, RED)
    tb(slide, f"⏱ {minutes} min", x + 0.06, y + 0.04, 1.22, 0.32,
       size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Wrap text for slide display ───────────────────────────────────────────────
def wrap_para(text, width=90):
    """Wrap a paragraph to ~width chars for slide display."""
    return "\n".join(textwrap.fill(text, width).split("\n"))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════

def course_title_slide(prs):
    s = new_slide(prs)
    bg(s, DARK)
    rect(s, 0, 0, 10, 0.09, RED)
    rect(s, 0, 7.41, 10, 0.09, RED)
    tb(s, COURSE_TITLE, 0.5, 1.3, 9, 1.6, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, "Teacher Slide Deck", 0.5, 2.95, 9, 0.7, size=22, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s, f"{GRADE}  ·  Semester Elective  ·  {SCHOOL_NAME}", 0.5, 3.75, 9, 0.5,
       size=13, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 2, 4.5, 6, 0.03, GOLD)
    tb(s, "Each lesson slide deck includes: reading slides · comprehension questions\n"
          "teaching content · activity directions · exit ticket\n"
          "Teacher notes (in Notes panel) include full talk tracks for every slide.",
       0.5, 4.7, 9, 1.0, size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    notes(s, "Welcome to the Warrior Prep Teacher Slide Deck.\n\n"
             "This deck is organized by unit and lesson. Each lesson follows the same structure:\n"
             "1. Lesson title & objectives\n2. Reading (text on screen)\n"
             "3. Comprehension questions\n4. Teaching content slides\n"
             "5. Activity slides (with the activity on screen for projection)\n"
             "6. Exit ticket\n\nAll teacher notes include talk tracks — suggested language, "
             "facilitation tips, and differentiation ideas.")


def unit_divider_slide(prs, unit):
    s = new_slide(prs)
    color = uc(unit)
    bg(s, color)
    rect(s, 0, 0, 10, 0.09, WHITE)
    rect(s, 0, 7.41, 10, 0.09, WHITE)
    tb(s, f"UNIT {unit['number']}", 0.5, 0.8, 9, 0.55, size=20, color=GOLD)
    tb(s, unit["name"], 0.5, 1.35, 9, 1.3, size=38, bold=True, color=WHITE)
    rect(s, 0.5, 2.75, 4.5, 0.05, GOLD)
    tb(s, unit["overview"], 0.5, 2.95, 9, 2.0, size=14, color=WHITE)

    # Vocab preview
    tb(s, "Unit Vocabulary", 0.5, 5.1, 9, 0.38, size=13, bold=True, color=GOLD)
    y = 5.55
    for v in unit["unit_vocabulary"][:3]:
        tb(s, f"• {v['word']}  —  {v['definition']}", 0.5, y, 9.3, 0.42, size=9.5, color=WHITE)
        y += 0.44

    notes(s,
          f"UNIT {unit['number']} OVERVIEW SLIDE\n\n"
          f"SAY: 'We're starting Unit {unit['number']}: {unit['name']}. Before we dive in, "
          f"let me read the overview with you.'\n\nRead the overview paragraph aloud. Then say: "
          f"'Let's look at the vocabulary we'll be building in this unit. Some of these words "
          f"you may have heard before — some may be new. Pay attention to the etymology column "
          f"on your notes pages. The word roots help you understand the meaning AND help you "
          f"decode other words with the same roots.'\n\n"
          f"PREVIEW QUESTION: Ask students to predict which of these vocabulary words they "
          f"think will be most challenging — and why. Quick pair-share.")


def vocab_deep_dive_slide(prs, unit):
    s = new_slide(prs)
    color = uc(unit)
    bg(s, LIGHT)
    header_bar(s, f"Unit {unit['number']} Vocabulary — Etymology & Meaning", color)
    y = 0.65
    for v in unit["unit_vocabulary"]:
        rect(s, 0.25, y, 9.5, 0.85, BLUE)
        tb(s, v["word"], 0.35, y + 0.04, 3.0, 0.38, size=13, bold=True, color=color)
        tb(s, f"⬅ {v['etymology']}", 0.35, y + 0.45, 4.5, 0.35, size=9, italic=True, color=GRAY)
        tb(s, v["definition"], 3.5, y + 0.1, 6.1, 0.65, size=11, color=DARK)
        y += 0.95

    notes(s,
          f"VOCABULARY DEEP DIVE\n\n"
          f"Walk through each word slowly. For each one:\n"
          f"1. Read the word aloud — have students repeat it.\n"
          f"2. Break down the etymology. Say: 'Look at the roots in this word. Does anyone "
          f"know another word with the same root?'\n"
          f"3. Read the definition. Ask: 'Put this in your own words.'\n\n"
          f"MORPHOLOGY TIP: Point out patterns across the unit — e.g., Latin roots that "
          f"appear in multiple words. Students who learn roots can decode thousands of "
          f"academic words they've never seen.\n\n"
          f"Have students write ALL words in their Cornell Notes vocabulary box now.")


def pretest_slide(prs, unit):
    s = new_slide(prs)
    color = uc(unit)
    bg(s, WHITE)
    header_bar(s, f"Unit {unit['number']} Pre-Test  —  What Do You Already Know?", color)
    sub_bar(s, "Answer each question honestly in your workbook. This is NOT graded for correctness.", color, 0.65, 0.35)
    y = 1.1
    for i, q in enumerate(unit["pretest"], 1):
        tb(s, f"{i}.  {q}", 0.3, y, 9.4, 0.52, size=12, color=DARK)
        y += 0.56
    notes(s,
          f"PRE-TEST — Unit {unit['number']}\n\n"
          f"SAY: 'Open your workbook to the Unit {unit['number']} pre-test page. Answer as "
          f"completely as you can. This is NOT graded for correctness — it is a snapshot of "
          f"where you start. Be honest — a blank answer is more useful than a guess.'\n\n"
          f"TIME: 8-10 minutes. Circulate but do not help with content.\n\n"
          f"AFTER: Collect or have students hold their pre-tests. You will return to them "
          f"at the end of the unit to compare. Quickly preview the unit: 'Here is what we "
          f"will be able to answer confidently by the end of this unit...'")


def bell_ringer_slide(prs, unit, lesson):
    """5-minute bell ringer that either connects to prior learning or launches the day."""
    color = uc(unit)
    label = lesson["lesson_label"]
    br = BELL_RINGERS.get(label, {})
    if not br:
        return

    s = new_slide(prs)
    bg(s, DARK)

    # Top accent stripe
    rect(s, 0, 0, 10, 0.12, color)
    rect(s, 0, 7.38, 10, 0.12, color)

    # Header
    rect(s, 0, 0.12, 10, 0.72, RED)
    tb(s, "🔔  BELL RINGER", 0.25, 0.18, 4.5, 0.52, size=20, bold=True, color=WHITE)
    tb(s, f"Lesson {label}  ·  5 minutes", 5.5, 0.25, 4.3, 0.38, size=13,
       color=GOLD, align=PP_ALIGN.RIGHT)

    # Bell ringer type chip
    type_colors = {"Launch": RED, "Review": color, "Anticipation": RGBColor(0x2E, 0x7D, 0x32)}
    chip_color = type_colors.get(br["type"], color)
    rect(s, 0.25, 0.95, 1.6, 0.36, chip_color)
    tb(s, br["type"].upper(), 0.3, 0.99, 1.5, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title
    tb(s, br["title"], 2.0, 0.95, 7.8, 0.44, size=17, bold=True, color=GOLD)

    # Prompt box
    rect(s, 0.25, 1.48, 9.5, 3.8, RGBColor(0x30, 0x30, 0x30))
    rect(s, 0.25, 1.48, 0.08, 3.8, chip_color)
    tb(s, br["prompt"], 0.45, 1.6, 9.1, 3.55, size=14, color=WHITE, line_spacing=20)

    # Share instruction
    rect(s, 0.25, 5.38, 9.5, 0.52, chip_color)
    tb(s, f"Share: {br['share']}", 0.38, 5.44, 9.2, 0.4, size=11, bold=True, color=WHITE)

    # Timer chip
    timer_chip(s, 5, x=8.4, y=6.1)

    # Reminder bar
    rect(s, 0.25, 6.05, 7.95, 0.45, RGBColor(0x38, 0x38, 0x38))
    tb(s, "📓  Write your response in your workbook before we share out.",
       0.38, 6.1, 7.7, 0.35, size=10, italic=True, color=GOLD)

    notes(s,
          f"BELL RINGER — Lesson {label}: {lesson['title']}\n\n"
          f"TYPE: {br['type']}\n\n"
          f"WHY THIS BELL RINGER: {br['connection']}\n\n"
          f"FACILITATION:\n"
          f"• Display this slide AS students enter. They should begin writing immediately — "
          f"no instruction needed.\n"
          f"• Give exactly 3 minutes of silent writing, then 1 minute of pair-share.\n"
          f"• Cold-call 1-2 students (30 sec each) before moving to the lesson title slide.\n"
          f"• Total target: 5 minutes from door to transition.\n\n"
          f"SHARE STRATEGY: {br['share']}\n\n"
          f"DO NOT OVER-DEBRIEF: The bell ringer is a warm engine, not the main event. "
          f"Acknowledge good responses, make the connection explicit, and move on.")


def lesson_title_slide(prs, unit, lesson):
    s = new_slide(prs)
    color = uc(unit)
    t = LESSON_TIMING.get(lesson["lesson_label"], {})
    bg(s, WHITE)
    rect(s, 0, 0, 10, 1.05, color)
    rect(s, 0, 0, 0.14, 7.5, color)
    tb(s, f"Lesson {lesson['lesson_label']}  |  Unit {unit['number']}: {unit['name']}",
       0.25, 0.06, 9.6, 0.38, size=12, color=GOLD)
    tb(s, lesson["title"], 0.25, 0.44, 9.5, 0.58, size=24, bold=True, color=WHITE)

    # Objectives
    sub_bar(s, "Learning Objectives", color, 1.12, 0.36)
    y = 1.55
    for obj in lesson["objectives"]:
        tb(s, f"◉  {obj}", 0.3, y, 5.6, 0.46, size=10.5, color=DARK)
        y += 0.48

    # Word focus
    rect(s, 6.2, 1.12, 3.65, 2.5, BLUE)
    tb(s, "Word Focus", 6.3, 1.14, 3.5, 0.33, size=10, bold=True, color=color)
    vf = lesson["vocab_focus"]
    tb(s, vf["word"], 6.3, 1.5, 3.5, 0.46, size=16, bold=True, color=color)
    tb(s, vf["etymology"], 6.3, 1.97, 3.5, 0.36, size=9, italic=True, color=GRAY)
    tb(s, vf["definition"], 6.3, 2.35, 3.5, 1.2, size=10.5, color=DARK)

    # 50-minute timing plan (visual bar)
    # Story discussion is folded into the reading/questions block for the bar
    # (reading + Qs + story ≈ same combined window; no separate timing entry)
    segments = [
        ("Bell\nRinger", t.get("bell", 5), RED),
        ("Intro &\nObjectives", t.get("title", 3), color),
        ("Reading &\nStory", t.get("read", 10), RGBColor(0x1A, 0x73, 0xE8)),
        ("Questions &\nDiscussion", t.get("qs", 7), RGBColor(0x2E, 0x7D, 0x32)),
        ("Content", t.get("content", 9), RGBColor(0x80, 0x00, 0x80)),
        ("Activity", t.get("activity", 11), RGBColor(0xE6, 0x5C, 0x00)),
        ("Exit\nTicket", t.get("exit", 5), DARK),
    ]
    total = sum(s2 for _, s2, _ in segments)

    rect(s, 0.25, 5.82, 9.5, 0.3, LIGHT)
    tb(s, f"50-Minute Lesson Plan", 0.25, 5.82, 4.0, 0.28, size=9, bold=True, color=DARK)
    tb(s, f"Total: {total} min", 7.0, 5.82, 2.5, 0.28, size=9, color=GRAY, align=PP_ALIGN.RIGHT)

    bar_x = 0.25
    bar_y = 6.15
    bar_h = 0.55
    bar_total_w = 9.5
    label_y = bar_y + bar_h + 0.04
    for seg_name, seg_min, seg_color in segments:
        seg_w = (seg_min / total) * bar_total_w
        rect(s, bar_x, bar_y, seg_w, bar_h, seg_color)
        if seg_w > 0.5:
            tb(s, str(seg_min), bar_x + 0.04, bar_y + 0.1, seg_w - 0.08, 0.35,
               size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, seg_name, bar_x, label_y, seg_w, 0.42,
           size=7, color=DARK, align=PP_ALIGN.CENTER)
        bar_x += seg_w

    scaffolding_labels = {
        1: "Level 1 — Heavy scaffold (most notes provided)",
        2: "Level 2 — Partial scaffold (key notes provided)",
        3: "Level 3 — Light scaffold (guiding questions only)",
        4: "Level 4 — Full independence (blank Cornell format)",
    }
    rect(s, 0, 7.2, 10, 0.3, RGBColor(0xE8, 0xE0, 0xD0))
    tb(s, f"Note-taking support: {scaffolding_labels.get(lesson['scaffolding'], '')}",
       0.3, 7.22, 9.4, 0.26, size=9, italic=True, color=DARK)

    notes(s,
          f"LESSON TITLE SLIDE — Lesson {lesson['lesson_label']}: {lesson['title']}\n\n"
          f"50-MINUTE BREAKDOWN:\n"
          f"  • Bell Ringer: {t.get('bell',5)} min\n"
          f"  • Intro & Objectives: {t.get('title',3)} min\n"
          f"  • Reading: {t.get('read',10)} min\n"
          f"  • Comprehension Questions: {t.get('qs',7)} min\n"
          f"  • Teaching Content: {t.get('content',9)} min\n"
          f"  • Activity: {t.get('activity',11)} min\n"
          f"  • Exit Ticket: {t.get('exit',5)} min\n\n"
          f"PACING NOTES:\n"
          f"• If the reading takes longer than planned, trim the content discussion — the "
          f"activity is the anchor for this lesson.\n"
          f"• If the activity is running long, give a 2-minute warning and have students "
          f"complete the most important section only — debrief with a partner.\n"
          f"• The exit ticket is non-negotiable. Build in a habit of ending with 5 clear "
          f"minutes. Students who know they must write at the end stay more engaged throughout.\n\n"
          f"SAY: 'Open your workbook to Lesson {lesson['lesson_label']}. Start with the "
          f"Cornell Notes page.'\n\n"
          f"WORD FOCUS: Read the word, etymology, and definition aloud. Ask: 'What other "
          f"words do you know that share these roots?' Give 30 seconds for pair-share.\n\n"
          f"OBJECTIVES: Read each one aloud. Say: 'By the end of class, you should be able "
          f"to do all three of these. Keep them in mind as we work through the lesson.'\n\n"
          f"SCAFFOLDING LEVEL {lesson['scaffolding']}: {scaffolding_labels.get(lesson['scaffolding'], '')}"
          f"\n\nTEACHER TIP: Early lessons have notes provided in the workbook. Gradually "
          f"reduce support across the semester so students build independence.")


def reading_slides(prs, unit, lesson):
    reading = READINGS.get(lesson["lesson_label"])
    if not reading:
        return
    color = uc(unit)

    # Reading title slide
    s = new_slide(prs)
    bg(s, color)
    rect(s, 0, 0, 10, 0.09, WHITE)
    rect(s, 0, 7.41, 10, 0.09, WHITE)
    tb(s, "READING", 0.5, 1.0, 9, 0.55, size=18, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s, reading["title"], 0.5, 1.6, 9, 1.4, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, f"Source: {reading['source']}", 0.5, 3.15, 9, 0.45, size=11, italic=True,
       color=GOLD, align=PP_ALIGN.CENTER)
    rect(s, 1.5, 3.8, 7, 0.04, GOLD)
    tb(s, "As you read: annotate using your symbols\n"
          "* = important   ?  = question   ! = surprising   → = connection   circle = vocab",
       0.5, 4.0, 9, 0.9, size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # Watch button — single linked shape (text run hyperlink = reliably clickable)
    yt = YOUTUBE.get(lesson["lesson_label"])
    if yt:
        link_btn(s, "▶  Watch Video", 0.4, 5.2, 2.4, 0.52,
                 GOLD, DARK, yt["url"], size=12)

    read_min = LESSON_TIMING.get(lesson["lesson_label"], {}).get("read", 10)
    timer_chip(s, read_min, x=8.5, y=6.5)

    notes(s,
          f"READING INTRO — '{reading['title']}'\n\n"
          f"SAY: 'Turn to the reading page in your workbook — it's right after the Cornell Notes "
          f"page. Today's reading is about {reading['title']}. Before you start, "
          f"make sure your annotation system is fresh in your mind.'\n\n"
          f"READING OPTIONS (choose based on your class):\n"
          f"  A) Independent silent reading (fastest)\n"
          f"  B) Popcorn read — students take turns reading paragraphs aloud\n"
          f"  C) Partner read — pairs read to each other\n"
          f"  D) You read it aloud while students follow and annotate\n\n"
          f"ANNOTATION REMINDER: Say 'I want to see at least 4 annotation marks per paragraph "
          f"by the time we finish. That means you are actually thinking, not just reading.'\n\n"
          f"TIME: 10 minutes for reading, 8 minutes for questions.")

    # Reading paragraphs — 2 paras per slide
    paragraphs = reading["text"]
    for i in range(0, len(paragraphs), 2):
        s = new_slide(prs)
        bg(s, WHITE)
        header_bar(s, f"Reading: {reading['title']}  (p. {i//2 + 1} of {(len(paragraphs)+1)//2})", color)
        y = 0.68
        for j, para in enumerate(paragraphs[i:i+2]):
            # Paragraph number
            rect(s, 0.2, y, 0.32, 0.32, color)
            tb(s, str(i + j + 1), 0.22, y + 0.02, 0.28, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            # Para text
            tb(s, wrap_para(para, 88), 0.6, y, 9.2, 2.65, size=11.5, color=DARK, line_spacing=16)
            # Annotation margin indicator
            rect(s, 9.7, y, 0.18, 2.5, LIGHT)
            tb(s, "mark\nhere", 9.71, y + 0.6, 0.16, 0.9, size=6, color=GRAY, align=PP_ALIGN.CENTER)
            y += 3.0

        notes(s,
              f"READING — Paragraph(s) {i+1}–{min(i+2, len(paragraphs))}\n\n"
              f"FACILITATION: Read paragraph by paragraph (or have students read). "
              f"After each paragraph, pause and ask: 'What symbol did you put in the margin and why?'\n\n"
              f"VOCABULARY SUPPORT: If students encounter an unknown word, direct them to:\n"
              f"  1. Read the surrounding sentences for context clues\n"
              f"  2. Break the word into parts (prefix, root, suffix)\n"
              f"  3. Make a best guess and keep reading\n\n"
              f"COMPREHENSION CHECK: After each paragraph, you can cold-call: "
              f"'In one sentence, what was that paragraph about?'")

    # Comprehension questions slide
    s = new_slide(prs)
    bg(s, LIGHT)
    header_bar(s, f"Reading Questions — '{reading['title']}'", color)
    sub_bar(s, "Answer in complete sentences. Use evidence from the reading. Work independently first, then discuss.", color, 0.65, 0.35)
    y = 1.1
    qs_min = LESSON_TIMING.get(lesson["lesson_label"], {}).get("qs", 7)
    timer_chip(s, qs_min)
    for i, q in enumerate(reading["questions"], 1):
        # Alternate row shading
        bg_clr = CREAM if i % 2 == 0 else WHITE
        rect(s, 0.2, y, 9.6, 0.54, bg_clr)
        tb(s, f"{i}.", 0.28, y + 0.08, 0.32, 0.38, size=12, bold=True, color=color)
        tb(s, wrap_para(q, 95), 0.65, y + 0.06, 9.0, 0.44, size=10.5, color=DARK)
        y += 0.57

    notes(s,
          f"COMPREHENSION QUESTIONS — '{reading['title']}'\n\n"
          f"SAY: 'Now turn to the questions section on your reading page. Work independently "
          f"for 6-8 minutes, then we will share out.'\n\n"
          f"DEBRIEF STRATEGY:\n"
          f"• Cold-call for Q1-2 (factual comprehension)\n"
          f"• Pair-share for Q3-4 (analysis/interpretation)\n"
          f"• Whole class discussion for Q5 (synthesis/personal connection)\n\n"
          f"DIFFERENTIATION:\n"
          f"• Students who finish early: write a follow-up question they would ask the author\n"
          f"• Students who struggle: work in pairs, or focus on Q1 and Q3 only\n\n"
          f"TRANSITION: After debrief, say: 'The reading gave us the foundation. Now let's "
          f"dig deeper into the content.'")


def story_slides(prs, unit, lesson):
    """Two slides per lesson: story text + discussion questions."""
    color = uc(unit)
    label = lesson["lesson_label"]
    s = STORIES.get(label)
    if not s:
        return

    # ── Slide 1: Story text ─────────────────────────────────────────────────
    sl = new_slide(prs)
    bg(sl, WHITE)
    header_bar(sl, f"Discussion Story: \"{s['title']}\"  —  Lesson {label}", color)
    rect(sl, 0, 0.58, 10, 0.3, DARK)
    tb(sl, "Characters:  Maya  ·  Jordan  ·  Priya  ·  Eli    |    "
          "Read together, then discuss.", 0.25, 0.61, 9.5, 0.24,
       size=10, italic=True, color=GOLD)

    # Story text — break into lines that fit the slide
    raw = s["story"].strip()
    paras = [p.strip() for p in raw.split("\n\n") if p.strip()]
    lines = []
    for para in paras:
        for ln in para.split("\n"):
            ln = ln.strip()
            if ln:
                lines.append(ln)

    y = 0.96
    for line in lines:
        if y > 7.05:
            break
        is_dialogue = line.startswith("<i>") or line.endswith("</i>")
        clean = line.replace("<i>", "").replace("</i>", "")
        font_size = 10.5 if is_dialogue else 10
        c = RGBColor(0x1A, 0x1A, 0x6E) if is_dialogue else DARK
        wrapped = wrap_para(clean, 94)
        line_count = wrapped.count("\n") + 1
        box_h = 0.22 * line_count
        tb(sl, wrapped, 0.25, y, 9.5, box_h + 0.05, size=font_size,
           color=c, italic=is_dialogue, line_spacing=15)
        y += box_h + 0.06

    notes(sl,
          f"DISCUSSION STORY — \"{s['title']}\"\n\n"
          f"SAY: 'Turn to the discussion story page in your workbook — it's right after "
          f"the informational reading. Today we're going to read a short story about four "
          f"8th graders: Maya, Jordan, Priya, and Eli. You'll be meeting them all semester.'\n\n"
          f"READ-ALOUD OPTIONS:\n"
          f"  A) Teacher reads aloud while students follow in workbook\n"
          f"  B) Assign a student to each character — they read that character's lines\n"
          f"  C) Students read silently in pairs, then discuss\n\n"
          f"AFTER READING (before the questions): 'Which character felt most familiar to "
          f"you — not who you want to be, but who you actually recognized?' "
          f"30-second pair-share before moving to questions.\n\n"
          f"TONE: These stories don't resolve neatly. That's intentional. The tension is "
          f"the discussion point. Resist the urge to tell students what the characters "
          f"should have done — let them figure it out.")

    # ── Slide 2: Discussion questions ──────────────────────────────────────
    sl2 = new_slide(prs)
    bg(sl2, CREAM)
    header_bar(sl2, f"Discussion Questions  —  \"{s['title']}\"", color)
    sub_bar(sl2, "Discuss with a partner. Then write your best thinking in your workbook.", color, 0.65, 0.35)
    timer_chip(sl2, 8)

    y2 = 1.1
    for i, q in enumerate(s["questions"], 1):
        bg_clr = WHITE if i % 2 == 1 else BLUE
        rect(sl2, 0.2, y2, 9.6, 0.68, bg_clr)
        rect(sl2, 0.2, y2, 0.42, 0.68, color)
        tb(sl2, str(i), 0.22, y2 + 0.14, 0.38, 0.38, size=13, bold=True,
           color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl2, wrap_para(q, 93), 0.7, y2 + 0.08, 9.0, 0.56, size=10.5, color=DARK)
        y2 += 0.72

    notes(sl2,
          f"DISCUSSION DEBRIEF\n\n"
          f"FACILITATION SEQUENCE:\n"
          f"1. Pairs discuss all 4 questions (5 min)\n"
          f"2. Cold-call Q1 — factual/observation level (who said/did what)\n"
          f"3. Pair-share Q3 or Q4 — personal connection (which character are you?)\n"
          f"4. Whole class for Q4 — the 'what would you do' question\n\n"
          f"TRANSITION: 'The story shows us the problem. The lesson content and activity "
          f"give us the tools. Let's go deeper.'\n\n"
          f"DISCUSSION NORMS REMINDER: Remind students that honest answers are more "
          f"valuable than impressive-sounding ones. The point is to recognize the real "
          f"situation — not to pretend they'd always make the right choice.\n\n"
          f"CONNECTION TO LESSON: Each question leads directly into the lesson's core "
          f"concept. If discussion stalls, use the question as a bridge: 'What we're "
          f"about to learn gives us the answer to that question.'")


def content_slides(prs, unit, lesson):
    color = uc(unit)
    for i, sec in enumerate(lesson["content_sections"]):
        s = new_slide(prs)
        bg(s, WHITE)
        header_bar(s, sec["heading"], color)

        # Body text
        tb(s, wrap_para(sec["body"], 88), 0.25, 0.68, 9.5, 3.8, size=13, color=DARK, line_spacing=18)

        # Discussion prompt after first section
        if i == 0:
            rect(s, 0.25, 4.65, 9.5, 1.1, GOLD)
            tb(s, "Think-Pair-Share:", 0.38, 4.7, 2.0, 0.38, size=11, bold=True, color=DARK)
            discussion = _get_discussion(lesson, i)
            tb(s, discussion, 0.38, 5.1, 9.2, 0.6, size=11, italic=True, color=DARK)
        elif i == 1:
            # Key takeaway box
            rect(s, 0.25, 4.65, 9.5, 1.1, BLUE)
            tb(s, "Key Idea to Note:", 0.38, 4.7, 2.2, 0.38, size=11, bold=True, color=color)
            tb(s, _get_key_idea(lesson, i), 0.38, 5.1, 9.2, 0.6, size=11, color=DARK)

        rect(s, 0, 7.18, 10, 0.32, LIGHT)
        tb(s, f"📝  Write key ideas in your Cornell Notes — RIGHT column   |   "
              f"Any questions? Add them to the LEFT column",
           0.25, 7.2, 9.5, 0.28, size=9, italic=True, color=GRAY)

        notes(s,
              f"CONTENT SLIDE: '{sec['heading']}'\n\n"
              f"TALK TRACK:\n{_get_talk_track(lesson, i, sec)}\n\n"
              f"AFTER READING: Ask students to add at least ONE item to their Cornell Notes "
              f"right column before you move on.\n\n"
              f"COMMON MISCONCEPTION TO ADDRESS: {_get_misconception(lesson, i)}\n\n"
              f"DIFFERENTIATION:\n"
              f"• Advanced: Ask students to identify a claim in the content that could be "
              f"questioned or that needs more evidence.\n"
              f"• Support: Provide the key sentence: 'The most important idea here is ___.'\n"
              f"  Have students copy it and then explain it in their own words.")

        # Activity slide (if this section has one)
        act = sec.get("activity")
        if act:
            activity_slide(prs, unit, lesson, sec, act, color)


def activity_slide(prs, unit, lesson, sec, act, color):
    at = act["type"]
    s = new_slide(prs)
    bg(s, LIGHT)
    header_bar(s, f"Activity: {_activity_title(at)}  —  Lesson {lesson['lesson_label']}", DARK)

    # Use lesson-specific activity time from the 50-min plan
    time_min = LESSON_TIMING.get(lesson["lesson_label"], {}).get("activity", _activity_time(at))
    timer_chip(s, time_min)

    directions = _activity_directions(at, act)
    tb(s, directions, 0.25, 0.65, 8.0, 1.2, size=11, italic=True, color=DARK)

    # Show the actual activity content on screen
    _render_activity_on_slide(s, at, act, color)

    notes(s,
          f"ACTIVITY: {_activity_title(at)}\n\n"
          f"TIME: {time_min} minutes\n\n"
          f"SETUP: {_activity_setup(at, act)}\n\n"
          f"FACILITATION:\n{_activity_facilitation(at, act)}\n\n"
          f"DEBRIEF: {_activity_debrief(at, act)}")


def exit_ticket_slide(prs, unit, lesson):
    color = uc(unit)
    et = lesson["exit_ticket"]
    et_type = et.get("type", "quick_write")

    s = new_slide(prs)
    bg(s, CREAM)
    header_bar(s, f"Exit Ticket  —  Lesson {lesson['lesson_label']}: {lesson['title']}", color)

    # Type label
    type_labels = {
        "quick_write": "Quick Write",
        "three_two_one": "3-2-1 Reflection",
        "reflection_scale": "Rating Scale + Reflection",
        "sentence_stems": "Sentence Stems",
        "vocabulary_check": "Vocabulary Check",
        "planning_response": "Planning Response",
        "postcard": "Postcard Write",
        "synthesis_write": "Synthesis Write",
    }
    tb(s, f"Format: {type_labels.get(et_type, et_type)}", 0.25, 0.65, 5.0, 0.35,
       size=11, italic=True, color=GRAY)
    timer_chip(s, 5)

    # Show actual prompt on screen
    y = 1.1
    if et_type in ("quick_write", "synthesis_write"):
        rect(s, 0.25, y, 9.5, 1.8, WHITE)
        tb(s, et["prompt"], 0.35, y + 0.08, 9.2, 1.65, size=12, color=DARK)
        y += 2.0
        _draw_lines_on_slide(s, y, 4)

    elif et_type == "three_two_one":
        for prompt in et["prompts"]:
            rect(s, 0.25, y, 9.5, 0.52, WHITE if y < 3 else LIGHT)
            tb(s, prompt, 0.35, y + 0.06, 9.2, 0.42, size=11, italic=True, color=DARK)
            y += 0.6
        _draw_lines_on_slide(s, y + 0.1, 2)

    elif et_type == "reflection_scale":
        rect(s, 0.25, y, 9.5, 1.1, WHITE)
        tb(s, et["prompt"], 0.35, y + 0.06, 9.2, 1.0, size=11, color=DARK)
        y += 1.2
        scale_labels = ["1 — Not at all", "2", "3 — Somewhat", "4", "5 — Yes!"]
        for i, label in enumerate(scale_labels):
            box_color = BLUE if i % 2 == 0 else WHITE
            rect(s, 0.25 + i * 1.85, y, 1.8, 0.56, box_color)
            tb(s, label, 0.3 + i * 1.85, y + 0.08, 1.7, 0.4,
               size=11, align=PP_ALIGN.CENTER, color=DARK)
        _draw_lines_on_slide(s, y + 0.8, 3)

    elif et_type == "sentence_stems":
        for prompt in et["prompts"]:
            rect(s, 0.25, y, 9.5, 0.45, WHITE)
            tb(s, prompt, 0.35, y + 0.05, 9.2, 0.38, size=11, italic=True, color=DARK)
            y += 0.55
            _draw_lines_on_slide(s, y, 2)
            y += 0.5

    elif et_type == "vocabulary_check":
        tb(s, et["prompt"], 0.3, y, 9.4, 0.6, size=11, color=DARK)
        y += 0.7
        for term, defn in zip(et["terms"], et["definitions"]):
            rect(s, 0.25, y, 3.0, 0.42, BLUE)
            rect(s, 3.3, y, 6.4, 0.42, WHITE)
            tb(s, term, 0.35, y + 0.05, 2.8, 0.32, size=11, bold=True, color=color)
            tb(s, defn, 3.4, y + 0.05, 6.2, 0.32, size=10, color=DARK)
            y += 0.48

    elif et_type == "planning_response":
        for prompt in et["prompts"]:
            rect(s, 0.25, y, 9.5, 0.45, WHITE)
            tb(s, prompt, 0.35, y + 0.05, 9.2, 0.38, size=11, italic=True, color=DARK)
            y += 0.6

    elif et_type == "postcard":
        rect(s, 0.25, y, 9.5, 1.3, WHITE)
        tb(s, et["prompt"], 0.35, y + 0.06, 9.2, 1.2, size=11, color=DARK)
        y += 1.45
        # Postcard template preview
        rect(s, 0.25, y, 5.0, 1.5, LIGHT)
        rect(s, 5.35, y, 4.4, 1.5, WHITE)
        tb(s, "[Write your postcard here]", 0.4, y + 0.5, 4.7, 0.5, size=10, italic=True, color=GRAY)
        tb(s, "To: My Future 9th-Grade Self\nTwo Rivers HS, Mendota Heights, MN 55118",
           5.45, y + 0.3, 4.2, 0.9, size=10, color=DARK)

    else:
        tb(s, str(et.get("prompt", "")), 0.3, y, 9.4, 2.0, size=12, color=DARK)
        _draw_lines_on_slide(s, y + 2.1, 3)

    notes(s,
          f"EXIT TICKET — Lesson {lesson['lesson_label']}\n\n"
          f"Format: {type_labels.get(et_type, et_type)}\n\n"
          f"SAY: 'You have 4-5 minutes for the exit ticket in your workbook. Complete it "
          f"before you leave — it's your ticket out the door.'\n\n"
          f"COLLECTION OPTIONS:\n"
          f"  A) Collect workbooks at the door (or remove and collect just the exit ticket page)\n"
          f"  B) Quick pair-share before collecting\n"
          f"  C) Cold-call 2 students to share their response with the class\n\n"
          f"HOW TO USE RESPONSES:\n"
          f"• Review quickly before next class — identify 2-3 strong examples to share "
          f"anonymously\n"
          f"• Note students who are confused — address in next lesson's warm-up\n"
          f"• If many students missed the key idea, start next lesson with a 2-minute "
          f"clarification before moving forward\n\n"
          f"NEXT LESSON: Connect the exit ticket insight to the opening of the next lesson.")


def posttest_slide(prs, unit):
    color = uc(unit)
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, f"Unit {unit['number']} Post-Test  —  How Far Have You Come?", color)
    tb(s, unit["posttest_prompt"], 0.3, 0.7, 9.4, 2.5, size=12.5, color=DARK)
    rect(s, 0.25, 3.35, 9.5, 0.72, BLUE)
    tb(s, "TIP: Your Cornell Notes are your best resource. Use the vocabulary, the key ideas "
          "from the readings, and the activities from this unit to write a thorough, specific response.",
       0.35, 3.44, 9.2, 0.55, size=11, italic=True, color=DARK)

    rect(s, 0.25, 4.25, 9.5, 0.55, GOLD)
    tb(s, "Growth Check: Compare your post-test to your pre-test. Circle 3 specific things "
          "you know now that you did not know at the start of this unit.",
       0.35, 4.32, 9.2, 0.42, size=11, bold=True, color=DARK)

    notes(s,
          f"POST-TEST — Unit {unit['number']}\n\n"
          f"SAY: 'Before you write, spend 2 minutes flipping through your Cornell Notes "
          f"and reading pages from this unit. This is an open-notes assessment.'\n\n"
          f"TIME: 12-15 minutes for writing, 3 minutes for growth comparison.\n\n"
          f"GROWTH CHECK: This is one of the most powerful moments in the unit. "
          f"Have students physically flip back to their pre-test and identify 3 specific "
          f"things that changed. Ask 2-3 students to share what changed for them.\n\n"
          f"ASSESSMENT: Score on: vocabulary use (correct, not just mentioned), "
          f"specificity (real examples, not vague claims), and connection to personal plan.\n\n"
          f"TRANSITION: Say 'Great work on Unit {unit['number']}. Let's take everything "
          f"we learned and carry it into Unit {unit['number']+1}.'")


# ── Activity rendering on slides ──────────────────────────────────────────────

def _draw_lines_on_slide(slide, y, n, color=RGBColor(0xC8, 0xD0, 0xDC)):
    for i in range(n):
        line_y = y + i * 0.45
        sh = slide.shapes.add_shape(1, Inches(0.25), Inches(line_y), Inches(9.5), Inches(0.02))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()


def _render_activity_on_slide(slide, at, act, color):
    """Put the actual activity content on the slide for projection."""
    y = 1.95

    if at == "rating_table":
        cols = act["columns"]
        # Header row
        rect(slide, 0.2, y, 5.5, 0.38, DARK)
        tb(slide, "Statement", 0.3, y + 0.05, 5.2, 0.28, size=10, bold=True, color=WHITE)
        for j, c in enumerate(cols):
            rect(slide, 5.75 + j * 0.7, y, 0.65, 0.38, color)
            tb(slide, c, 5.77 + j * 0.7, y + 0.05, 0.62, 0.28, size=10, bold=True,
               color=WHITE, align=PP_ALIGN.CENTER)
        y += 0.42
        for i, row in enumerate(act["rows"]):
            bg_clr = CREAM if i % 2 == 0 else WHITE
            rect(slide, 0.2, y, 5.5, 0.44, bg_clr)
            tb(slide, row, 0.28, y + 0.06, 5.3, 0.34, size=9.5, color=DARK)
            for j in range(len(cols)):
                rect(slide, 5.75 + j * 0.7, y, 0.65, 0.44, LIGHT)
                tb(slide, "☐", 5.8 + j * 0.7, y + 0.04, 0.55, 0.36, size=14,
                   color=DARK, align=PP_ALIGN.CENTER)
            y += 0.46
            if y > 7.1:
                break

    elif at == "audit_table":
        cols = act["columns"]
        rect(slide, 0.2, y, 3.5, 0.38, DARK)
        tb(slide, "Distractor", 0.3, y + 0.05, 3.3, 0.28, size=10, bold=True, color=WHITE)
        for j, c in enumerate(cols):
            rect(slide, 3.75 + j * 1.45, y, 1.4, 0.38, color)
            tb(slide, c, 3.77 + j * 1.45, y + 0.05, 1.36, 0.28, size=9, bold=True,
               color=WHITE, align=PP_ALIGN.CENTER)
        y += 0.42
        for i, row in enumerate(act["rows"]):
            bg_clr = CREAM if i % 2 == 0 else WHITE
            rect(slide, 0.2, y, 3.5, 0.44, bg_clr)
            tb(slide, row, 0.28, y + 0.06, 3.3, 0.34, size=9.5, color=DARK)
            for j in range(len(cols)):
                rect(slide, 3.75 + j * 1.45, y, 1.4, 0.44, LIGHT)
                tb(slide, "☐", 4.1 + j * 1.45, y + 0.04, 0.7, 0.36, size=14,
                   color=DARK, align=PP_ALIGN.CENTER)
            y += 0.46
            if y > 7.1:
                break

    elif at == "matrix_sort":
        # Draw the Eisenhower matrix
        labels = [("URGENT + IMPORTANT\n→ DO IT NOW", RED, DARK),
                  ("NOT URGENT + IMPORTANT\n→ SCHEDULE IT", color, DARK),
                  ("URGENT + NOT IMPORTANT\n→ MINIMIZE IT", RGBColor(0xAA, 0xAA, 0xAA), DARK),
                  ("NOT URGENT + NOT IMPORTANT\n→ ELIMINATE IT", LIGHT, DARK)]
        positions = [(0.2, y), (5.1, y), (0.2, y + 2.2), (5.1, y + 2.2)]
        for (lbl, fill, tc), (px, py) in zip(labels, positions):
            rect(slide, px, py, 4.7, 2.1, fill)
            tb(slide, lbl, px + 0.1, py + 0.15, 4.5, 0.8, size=11, bold=True, color=tc)
        # Task list
        tb(slide, "Tasks to sort:", 0.2, y + 4.6, 9.6, 0.3, size=10, bold=True, color=DARK)
        task_text = "  ".join([f"{i}. {t}" for i, t in enumerate(act["tasks"][:4], 1)])
        tb(slide, task_text, 0.2, y + 4.95, 9.6, 0.38, size=9.5, color=DARK)
        task_text2 = "  ".join([f"{i}. {t}" for i, t in enumerate(act["tasks"][4:], 5)])
        tb(slide, task_text2, 0.2, y + 5.35, 9.6, 0.38, size=9.5, color=DARK)

    elif at == "interest_table":
        hdrs = ["Career Field", "Description", "My Interest (1–5)", "Career I Know"]
        widths = [1.6, 2.5, 1.4, 2.7]
        rect(slide, 0.2, y, 9.6, 0.38, DARK)
        x = 0.25
        for hdr, w in zip(hdrs, widths):
            tb(slide, hdr, x, y + 0.05, w, 0.28, size=9, bold=True, color=WHITE)
            x += w + 0.08
        y += 0.42
        for i, (name, desc) in enumerate(act["rows"][:5]):
            bg_clr = CREAM if i % 2 == 0 else WHITE
            rect(slide, 0.2, y, 9.6, 0.52, bg_clr)
            x = 0.25
            for val, w in zip([name, desc, "", ""], widths):
                tb(slide, val, x, y + 0.06, w, 0.4, size=9, color=DARK)
                x += w + 0.08
            y += 0.54
            if y > 7.0:
                break

    elif at == "venn_diagram":
        # Draw two overlapping circles (approximated with rectangles + label)
        rect(slide, 0.25, y, 4.0, 3.8, BLUE)
        rect(slide, 5.75, y, 4.0, 3.8, GOLD)
        rect(slide, 3.5, y + 0.3, 3.0, 3.2, LIGHT)
        tb(slide, act["label_a"], 0.35, y + 0.1, 3.8, 0.42, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        tb(slide, act["label_b"], 5.85, y + 0.1, 3.8, 0.42, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        tb(slide, "Shared / Similar", 3.55, y + 0.4, 2.9, 0.38, size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        tb(slide, "(Differences)", 0.5, y + 1.5, 2.8, 0.38, size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
        tb(slide, "(Differences)", 6.0, y + 1.5, 2.8, 0.38, size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    elif at == "smart_goal_template":
        labels_letters = ["S — Specific", "M — Measurable", "A — Achievable", "R — Relevant", "T — Time-bound"]
        hints = [f[1] for f in act["fields"]]
        for i, (lbl, hint) in enumerate(zip(labels_letters, hints)):
            bg_clr = CREAM if i % 2 == 0 else WHITE
            rect(slide, 0.2, y, 9.6, 0.68, bg_clr)
            rect(slide, 0.2, y, 1.6, 0.68, color)
            tb(slide, lbl, 0.28, y + 0.12, 1.46, 0.44, size=11, bold=True, color=WHITE)
            tb(slide, hint, 1.88, y + 0.16, 7.8, 0.36, size=10, italic=True, color=GRAY)
            _draw_lines_on_slide(slide, y + 0.72, 1)
            y += 0.72
            if y > 6.8:
                break

    elif at == "gpa_calculator":
        hdrs = ["Course", "Grade", "Credits", "GPA Points", "Weighted"]
        widths = [2.8, 1.0, 1.0, 1.8, 1.8]
        rect(slide, 0.2, y, 9.6, 0.38, DARK)
        x = 0.25
        for hdr, w in zip(hdrs, widths):
            tb(slide, hdr, x, y + 0.05, w, 0.28, size=10, bold=True, color=WHITE)
            x += w + 0.05
        y += 0.42
        for i, c in enumerate(act["courses"]):
            bg_clr = CREAM if i % 2 == 0 else WHITE
            rect(slide, 0.2, y, 9.6, 0.44, bg_clr)
            x = 0.25
            for val, w in zip([c[0], c[1], c[2], "", ""], widths):
                tb(slide, val, x, y + 0.07, w, 0.3, size=9.5, color=DARK)
                x += w + 0.05
            y += 0.46
        rect(slide, 0.2, y, 9.6, 0.44, GOLD)
        tb(slide, "TOTAL GPA →", 0.3, y + 0.07, 3.0, 0.3, size=10, bold=True, color=DARK)

    elif at == "time_map":
        cols = act["columns"]
        col_w = 4.0
        rect(slide, 0.2, y, 1.5, 0.38, DARK)
        tb(slide, "Time", 0.28, y + 0.05, 1.35, 0.28, size=10, bold=True, color=WHITE)
        for j, c in enumerate(cols):
            rect(slide, 1.75 + j * (col_w + 0.15), y, col_w, 0.38, color)
            tb(slide, c, 1.78 + j * (col_w + 0.15), y + 0.05, col_w - 0.06, 0.28,
               size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += 0.42
        sample_hours = list(range(7, 23, 2))  # 7am to 10pm every 2 hours
        for i, h in enumerate(sample_hours[:8]):
            bg_clr = CREAM if i % 2 == 0 else WHITE
            label = f"{h % 12 or 12}:00 {'AM' if h < 12 else 'PM'}"
            rect(slide, 0.2, y, 1.5, 0.52, bg_clr)
            tb(slide, label, 0.28, y + 0.1, 1.36, 0.32, size=10, color=DARK, align=PP_ALIGN.CENTER)
            for j in range(len(cols)):
                rect(slide, 1.75 + j * (col_w + 0.15), y, col_w, 0.52, WHITE)
            y += 0.54
            if y > 7.0:
                break

    elif at == "final_schedule_with_rationale":
        hdrs = ["Period", "Course Choice", "Why I chose it", "Alternate"]
        widths = [0.7, 1.9, 3.5, 2.3]
        rect(slide, 0.2, y, 9.6, 0.38, color)
        x = 0.25
        for hdr, w in zip(hdrs, widths):
            tb(slide, hdr, x, y + 0.05, w, 0.28, size=10, bold=True, color=WHITE)
            x += w + 0.1
        y += 0.42
        for i in range(1, min(act["periods"] + 1, 8)):
            bg_clr = CREAM if i % 2 == 0 else WHITE
            rect(slide, 0.2, y, 9.6, 0.56, bg_clr)
            rect(slide, 0.2, y, 0.7, 0.56, LIGHT)
            tb(slide, f"P{i}", 0.25, y + 0.1, 0.58, 0.36, size=11, bold=True,
               color=DARK, align=PP_ALIGN.CENTER)
            y += 0.58
            if y > 7.0:
                break

    elif at in ("question_prep",):
        for i in range(1, min(act["num_questions"] + 1, 6)):
            rect(slide, 0.2, y, 9.6, 0.68, CREAM if i % 2 == 0 else WHITE)
            rect(slide, 0.2, y, 0.55, 0.68, color)
            tb(slide, str(i), 0.22, y + 0.12, 0.5, 0.44, size=14, bold=True,
               color=WHITE, align=PP_ALIGN.CENTER)
            tb(slide, "My question:", 0.85, y + 0.08, 2.5, 0.3, size=10, italic=True, color=GRAY)
            _draw_lines_on_slide(slide, y + 0.44, 1)
            y += 0.74
            if y > 7.0:
                break

    elif at == "activities_checklist":
        # Show top 3 categories
        cats = list(act["categories"].items())[:3]
        for cat, items in cats:
            rect(slide, 0.2, y, 9.6, 0.35, color)
            tb(slide, cat, 0.3, y + 0.05, 9.3, 0.25, size=11, bold=True, color=WHITE)
            y += 0.38
            trios = [items[i:i+3] for i in range(0, min(len(items), 9), 3)]
            for trio in trios[:2]:
                for j, item in enumerate(trio):
                    rect(slide, 0.2 + j * 3.2, y, 3.1, 0.4, CREAM if j % 2 == 0 else WHITE)
                    tb(slide, f"☐  {item}", 0.3 + j * 3.2, y + 0.06, 2.95, 0.3, size=10, color=DARK)
                y += 0.43
            if y > 6.8:
                break

    elif at in ("timed_write",):
        rect(slide, 0.2, y, 9.6, 0.7, RED)
        tb(slide, f"⏱ TIMED WRITE — {act.get('minutes', 12)} MINUTES", 0.3, y + 0.06,
           9.3, 0.58, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += 0.8
        rect(slide, 0.2, y, 9.6, 0.6, BLUE)
        tb(slide, "PROMPT: What is the most important skill a student needs to succeed in "
                  "high school? Use specific evidence from your own experience or this course.",
           0.3, y + 0.08, 9.3, 0.5, size=12, color=DARK)
        y += 0.75
        tb(slide, "STEP 1: Plan (2 min) — jot your main idea and 2-3 points below:", 0.25, y, 9.0, 0.38, size=11, bold=True, color=DARK)
        y += 0.45
        _draw_lines_on_slide(slide, y, 2)
        y += 1.0
        tb(slide, "STEP 2: DRAFT — write without stopping. Don't edit. Just go.", 0.25, y, 9.0, 0.38, size=11, bold=True, color=DARK)

    elif at in ("role_play",):
        tb(slide, "Self-Advocacy Script:", 0.25, y, 9.0, 0.38, size=12, bold=True, color=color)
        y += 0.45
        rect(slide, 0.25, y, 9.5, 1.1, BLUE)
        tb(slide, '"Hi [teacher\'s name], I was working on [assignment] and I\'m struggling with '
                  '[specific part]. I already tried [what I did]. Could you help me understand [question]?"',
           0.35, y + 0.1, 9.2, 0.9, size=13, italic=True, color=DARK)
        y += 1.25
        for i, sc in enumerate(act["scenarios"][:3]):
            rect(slide, 0.25, y, 9.5, 0.55, CREAM if i % 2 == 0 else WHITE)
            tb(slide, f"Scenario {i+1}:", 0.35, y + 0.06, 1.5, 0.38, size=10, bold=True, color=color)
            tb(slide, sc, 1.9, y + 0.1, 7.7, 0.38, size=10.5, color=DARK)
            y += 0.6

    elif at == "backward_plan":
        rect(slide, 0.2, y, 9.6, 0.38, DARK)
        for x, hdr in [(0.28, "Week"), (1.18, "Tasks to Complete by This Date"), (7.45, "Mini-Deadline")]:
            tb(slide, hdr, x, y + 0.05, 2.0, 0.28, size=10, bold=True, color=WHITE)
        y += 0.42
        for w in range(act["weeks"], 0, -1):
            bg_clr = CREAM if w % 2 == 0 else WHITE
            rect(slide, 0.2, y, 9.6, 0.62, bg_clr)
            rect(slide, 0.2, y, 0.88, 0.62, color)
            tb(slide, f"Wk {w}", 0.24, y + 0.12, 0.78, 0.38, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _draw_lines_on_slide(slide, y + 0.44, 1)
            y += 0.64
            if y > 7.0:
                break

    elif at == "study_plan_template":
        for i, field in enumerate(act["fields"][:6]):
            rect(slide, 0.2, y, 9.6, 0.6, CREAM if i % 2 == 0 else WHITE)
            tb(slide, field, 0.3, y + 0.1, 9.3, 0.42, size=11, bold=True, color=color)
            y += 0.65
            if y > 7.0:
                break

    else:
        # Generic: show directions as text
        tb(slide, _activity_directions(at, act), 0.25, y, 9.5, 3.5, size=12, color=DARK)


# ── Helper functions for slide content ───────────────────────────────────────

def _activity_title(at):
    titles = {
        "rating_table": "EF Self-Assessment",
        "audit_table": "Distraction Audit",
        "matrix_sort": "Eisenhower Matrix",
        "interest_table": "Career Field Inventory",
        "two_column_chart": "Course Connection Chart",
        "venn_diagram": "Compare Two Pathways",
        "guided_reflection": "Career Pathway Reflection",
        "smart_goal_template": "Write Your SMART Goal",
        "study_plan_template": "Build a Study Plan",
        "backward_plan": "Backward Plan a Project",
        "role_play": "Self-Advocacy Role Play",
        "annotation_practice": "Annotate the Passage",
        "timed_write": "Timed Write",
        "checklist": "Revision Checklist",
        "gpa_calculator": "GPA Calculator Practice",
        "habits_checklist": "Habits That Protect Your GPA",
        "commitment_card": "The Try One Thing Commitment",
        "activities_checklist": "TRHS Activity Exploration",
        "time_map": "Daily Time Map",
        "schedule_grid": "9th Grade Schedule Preview",
        "full_schedule_builder": "Build Your Mock Schedule",
        "year_plan_summary": "My 9th Grade Year Plan",
        "final_schedule_with_rationale": "Finalize Your Schedule",
        "question_prep": "Panel Question Preparation",
        "panel_notes": "Student Panel Notes",
    }
    return titles.get(at, at.replace("_", " ").title())


def _activity_time(at):
    times = {
        "rating_table": 5, "audit_table": 7, "matrix_sort": 10,
        "interest_table": 8, "two_column_chart": 8, "venn_diagram": 10,
        "guided_reflection": 12, "smart_goal_template": 12, "study_plan_template": 10,
        "backward_plan": 10, "role_play": 10, "annotation_practice": 12,
        "timed_write": 14, "checklist": 5, "gpa_calculator": 8,
        "habits_checklist": 4, "commitment_card": 5, "activities_checklist": 8,
        "time_map": 12, "schedule_grid": 10, "full_schedule_builder": 15,
        "year_plan_summary": 12, "final_schedule_with_rationale": 12,
        "question_prep": 8, "panel_notes": 25,
    }
    return times.get(at, 8)


def _activity_directions(at, act):
    dirs = {
        "rating_table": "Rate yourself honestly on each statement (1–5). This is for YOUR insight, not a grade. Total your scores and identify your top strength and biggest growth area.",
        "audit_table": "Circle how often each distractor affects your study time. Then: pick your #1 distractor and write a specific plan to address it in the box below the table.",
        "matrix_sort": "Draw the Eisenhower Matrix in your workbook (4 quadrants). Place each numbered task in the correct box. Then: what should you do FIRST? Share with a partner.",
        "interest_table": "Rate your interest in each career field (1–5). List one career you already know from each field. We'll share out after — be ready to explain your top rating.",
        "two_column_chart": "Choose the TWO courses from today's reading that interest you most. Use the Academic Planning Guide to complete the chart for each one.",
        "venn_diagram": "Choose two career pathways we've discussed. In the outer sections, list what's different about each. In the overlap, list what they have in common.",
        "guided_reflection": "Complete all five prompts using complete sentences. Use specific course names and pathway titles from the Academic Planning Guide.",
        "smart_goal_template": "Fill in each section of the SMART framework. Your goal must be school-related and must meet ALL five criteria. We'll share with a partner at the end.",
        "study_plan_template": "Identify a real upcoming test or major assignment. Build a spaced study plan with at least 3 sessions. Be specific: what will you DO in each session?",
        "backward_plan": "Start at your FINAL deadline. Work backward: what needs to be done each week? Give each task its own mini-deadline. Enter these in your planner.",
        "role_play": "Practice the self-advocacy script with a partner. Switch roles for each scenario. Listener: give feedback — was the request specific? Did they explain what they tried?",
        "annotation_practice": "Apply your annotation symbols to the passage. Aim for 4+ marks per paragraph. Then write your 2-sentence summary WITHOUT looking back at the passage.",
        "timed_write": "STEP 1: Plan for 2 minutes only. STEP 2: Draft for 10 minutes — do NOT stop, do NOT edit. STEP 3: Exchange with a partner for the revision checklist.",
        "checklist": "Exchange drafts with a partner. Use the checklist to give SPECIFIC feedback — not just 'good' or 'needs work.' Mark each item yes or no and explain.",
        "gpa_calculator": "Calculate the sample student's GPA using the TRHS grading scale. Then: which ONE grade improvement would have the biggest impact on their GPA? Show your work.",
        "habits_checklist": "Check the habits you already have. Circle the ones you want to build. Put a star next to the ONE habit you will commit to starting this week.",
        "commitment_card": "Fill in each field honestly and completely. This is a real commitment — write something you actually mean. You will revisit this at the end of the unit.",
        "activities_checklist": "Circle every activity that interests you — even a little. Put a star next to your TOP THREE. Be honest, not impressive.",
        "time_map": "Fill in a REALISTIC schedule — not an ideal one. Include sleep, meals, and commute. Then: is this sustainable? What would you change?",
        "schedule_grid": "Use the graduation plan (pages 12-13 of your Academic Planning Guide) to fill in required 9th grade courses. Then choose electives for your open periods.",
        "full_schedule_builder": "Build your complete 7-period mock schedule for BOTH semesters. Check prerequisites for every elective. Connect electives to your Unit 2 career interests.",
        "year_plan_summary": "Complete all six fields using your work from Lessons 4.1-4.4. This is the central document of your 9th Grade Year Plan.",
        "final_schedule_with_rationale": "Write your final choices and a one-sentence rationale for each. Identify one alternate for every elective. Make each choice deliberate.",
        "question_prep": "Write 5 genuine questions for the student panel. Make them specific and story-inviting. Avoid yes/no questions. Think about what you REALLY want to know.",
        "panel_notes": "Listen actively. Write the speaker's name and key ideas. Circle what surprises you. Star what you want to remember. Save your most important insight for the exit ticket.",
    }
    return dirs.get(at, "Complete the activity in your workbook following the instructions on the page.")


def _activity_setup(at, act):
    setups = {
        "matrix_sort": "Students will need to draw the matrix in their workbook (a large open area is provided). Consider drawing it on the board first so students can see the structure.",
        "timed_write": "Announce the timer clearly. At 2 minutes, say 'Planning time is up — begin your draft.' At 12 minutes total, say 'Finish your current sentence and stop.'",
        "role_play": "Have students find a partner before you give directions. Model the script once with a student volunteer BEFORE pairs begin practicing.",
        "panel_notes": "Brief panelists privately before class. Remind them: be honest, not just positive. Ask them to prepare one moment they struggled and how they handled it.",
        "gpa_calculator": "Write the TRHS grading scale on the board: A=4.0, A-=3.67, B+=3.33, B=3.0, etc. Students will need it to calculate GPA points.",
    }
    return setups.get(at, "Read the directions aloud. Clarify before students begin. Circulate during the activity.")


def _activity_facilitation(at, act):
    tips = {
        "rating_table": "Walk the room — encourage honest ratings, not inflated ones. After 4 minutes, cold-call: 'Who rated themselves a 2 or lower on one of these? What does that area look like for you?'",
        "audit_table": "Some students may not want to admit to social media distraction — normalize it. Say: 'Most adults have this problem too. The point is to build awareness, not shame.'",
        "matrix_sort": "Circulate during sorting. Common mistakes: students put everything in 'Urgent + Important.' Help them distinguish urgency from importance.",
        "smart_goal_template": "The most common problem: vague goals. If you see 'I want to be more organized,' push back: 'What does organized look like? How will you measure it? By when?'",
        "timed_write": "Walk the room. Students who freeze: whisper 'Just write your first thought — don't judge it.' Students who finish early: tell them to add one more piece of evidence.",
        "role_play": "After each scenario, give 30 seconds of feedback before switching. Key feedback: Was it specific? Did they say what they tried? Did they ask a clear question?",
        "gpa_calculator": "After calculations, ask: 'Which course would have the most impact if the grade went up one letter? Why?' This builds understanding of weighted grades.",
        "full_schedule_builder": "Circulate and check prerequisites. Many students will not have checked these. A common mistake: wanting to take advanced courses without the prerequisite.",
        "panel_notes": "During the panel: intervene if panelists give vague answers. Say: 'Can you give a specific example of that?' Let student questions drive the discussion.",
    }
    return tips.get(at, "Circulate during the activity. Ask guiding questions. Call time and debrief with 2-3 student shares.")


def _activity_debrief(at, act):
    debriefs = {
        "rating_table": "Ask: 'What was your lowest rating? What does that tell you about where to focus first in this unit?'",
        "matrix_sort": "Compare matrices: Do partners agree on every task? Where did you disagree — and why? Connect back to the reading: what does this tell you about how you currently use your time?",
        "interest_table": "Quick poll: How many gave a 5 to Engineering? Health Sciences? Business? etc. Ask one student from each field to share one career they listed.",
        "smart_goal_template": "Partners share goals. The listener's job: identify which SMART elements are strong and which need more specificity.",
        "timed_write": "Have 2 volunteers share one sentence they're proud of. Then: what was hard about writing without stopping? How does that connect to the idea of writing stamina?",
        "gpa_calculator": "Compare answers. Ask: 'If this were YOUR transcript, what would you do differently? What habit would you change starting today?'",
        "full_schedule_builder": "Have students share their most interesting elective choice and why. Connect to Unit 2: 'How does this connect to the career pathway you identified?'",
        "panel_notes": "Give 3 minutes for students to write their exit ticket before the panel ends. The most important thing you heard becomes the exit ticket.",
    }
    return debriefs.get(at, "Cold-call 2-3 students to share. Connect back to the lesson objective.")


def _get_discussion(lesson, i):
    prompts = {
        "1.1": "Think about your day so far. When did you need to use executive functioning skills? Which ones?",
        "1.2": "Is there a task you've been putting off because it's 'not due yet'? Is it urgent? Is it important?",
        "1.3": "How many tasks do you have open on your phone or computer right now? What does that cost you?",
        "1.4": "Think of a goal you set for yourself in the past. Was it SMART? What happened to it?",
        "2.1": "Before today, did you know what 'career pathways' meant? How might knowing this change how you think about school?",
        "2.2": "Would you rather take a traditional course or a CAPS internship experience? What does your answer tell you about how you learn?",
        "2.3": "Does your career interest fall more in Arts, Business, or Human Services? Or somewhere else?",
        "2.4": "What career field do you still know the least about? What would help you learn more?",
        "3.1": "When was the last time you read something difficult and actually remembered it? What made the difference?",
        "3.2": "Think of the best piece of writing you've done. What made it good? Did you revise it?",
        "3.3": "Honestly: how do you usually study before a big test? Does it work? How do you know?",
        "3.4": "Think of a time you hit a wall academically. What did you do? What would you do differently now?",
        "4.1": "Did you know Two Rivers had these many graduation requirements? What surprises you?",
        "4.2": "What 9th grade course are you most looking forward to? Most nervous about?",
        "4.3": "What has stopped you from trying something new in the past? Is that still a good reason?",
        "4.4": "How many hours of sleep do you typically get on a school night? How does it affect you?",
        "5.1": "Have you ever checked your grades in a portal? How often? What stopped you from doing it more?",
        "5.2": "What do you know about how 8th grade is different from what you expect in 9th grade?",
        "5.3": "If you had to submit your course registration TODAY, what would you put down? How confident do you feel?",
        "5.4": "What is one thing you genuinely want to know from a current 9th grader — something you haven't heard anyone address yet?",
    }
    return prompts.get(lesson["lesson_label"], "Turn to a partner: what was the most important idea from that section?")


def _get_key_idea(lesson, i):
    ideas = {
        "1.1": "EF skills are not fixed — they are learnable. Every strategy you practice in this course is strengthening real neural pathways.",
        "1.2": "The most important tasks are often the ones that are NOT urgent yet. High performers schedule these before they become crises.",
        "1.3": "A phone on your desk — even face-down — reduces cognitive performance. Environment design beats willpower every time.",
        "1.4": "Written goals are significantly more likely to be achieved than goals you just keep in your head. Writing = commitment.",
        "2.1": "Every course you take at Two Rivers connects to at least one of the six Minnesota career fields. Your electives are not random.",
        "2.2": "CAPS programs require an application and a strong attendance record — because real workplaces demand exactly that.",
        "2.3": "The CAPS Computer Science program at Two Rivers earns students an industry-recognized A+ certification alongside their high school credits.",
        "2.4": "Career planning at 13-14 is not about deciding your life — it's about exploring with purpose.",
        "3.1": "Annotation is not decoration — it is evidence that your brain is actively processing the text.",
        "3.2": "The first draft's only job is to exist. Clarity comes in revision, not in drafting.",
        "3.3": "Struggle during retrieval practice is not a sign of failure — it is the engine of durable learning.",
        "3.4": "Self-advocacy is a skill, and like all skills, it improves with deliberate practice.",
        "4.1": "Prerequisite chains matter: choices you make in 9th grade determine which courses are available to you in 10th, 11th, and 12th.",
        "4.2": "Your 9th grade registration is final at Two Rivers. The time to plan carefully is now, not in December.",
        "4.3": "Students with extracurricular involvement have higher GPAs — not because they are smarter, but because structured activities create discipline.",
        "4.4": "Protect your homework time FIRST. High performers decide in advance when homework happens and treat it as non-negotiable.",
        "5.1": "Infinite Campus shows your grade in real time. Students who monitor their grades regularly catch problems faster and recover faster.",
        "5.2": "BARR amplifies the effort you bring. Students who show up and communicate benefit enormously; students who disengage miss most of it.",
        "5.3": "Alternates are not an afterthought — treat them with the same intentionality as your first choices.",
        "5.4": "The best questions invite stories, not opinions. 'Tell me about a time when...' almost always produces more useful answers.",
    }
    return ideas.get(lesson["lesson_label"], "The most important takeaway from today connects directly to your SMART goal.")


def _get_talk_track(lesson, i, sec):
    label = lesson["lesson_label"]
    if i == 0:
        return (f"SAY: 'Let's dig into the content. The reading gave us the big picture — "
                f"now we are going to break down the details.'\n\n"
                f"Read the section aloud or paraphrase it. Stop after each bold term: "
                f"'What does that mean in your own words?' Give 20 seconds before continuing.\n\n"
                f"Direct students to their Cornell Notes right column: 'Write down at least "
                f"3 key ideas from this section before we move on.'")
    else:
        return (f"SAY: 'We are going deeper now. This section builds directly on what we "
                f"just discussed.'\n\nBefore reading: ask students to predict — "
                f"'Based on what you know so far, what do you think this section will cover?'\n\n"
                f"After reading: 'What is the connection between this section and the first one?'")


def _get_misconception(lesson, i):
    misconceptions = {
        "1.1": "Students often think EF is about intelligence — it's not. It's about habits and strategies.",
        "1.2": "Students often confuse 'urgent' and 'important.' Many urgent things are not actually important.",
        "1.3": "Many students believe they CAN multitask effectively. The research is clear: they cannot.",
        "1.4": "Students often think willpower is the key to self-regulation. The evidence shows strategies matter far more.",
        "2.1": "Students sometimes think career fields are rigid boxes — they are organizing tools, not limitations.",
        "2.2": "Students assume CAPS programs are only for students who already know exactly what they want to do.",
        "2.3": "Students often underestimate the business field — it includes entrepreneurship, not just corporate work.",
        "2.4": "Students sometimes think planning careers in 8th grade means making permanent decisions — it doesn't.",
        "3.1": "Students believe highlighting equals learning. Research shows it is one of the LEAST effective strategies.",
        "3.2": "Students think editing while drafting saves time — it actually doubles the time and fragments thinking.",
        "3.3": "Students equate familiarity with mastery. They feel like they 'know it' after re-reading, but can't retrieve it.",
        "3.4": "Students think asking for help shows weakness — it actually shows metacognitive awareness and earns teacher respect.",
        "4.1": "Students sometimes think electives don't matter for graduation — they do: 5.5 credits worth.",
        "4.2": "Students assume Warrior Seminar is an easy filler class — it is actually one of the most directly useful courses in 9th grade.",
        "4.3": "Students assume activities hurt their GPA — research shows the opposite.",
        "4.4": "Students believe sleep is the first thing to sacrifice when busy. It is actually the last thing you can afford to sacrifice.",
        "5.1": "Students think GPA is set in stone after a bad semester — it can absolutely be recovered, but it takes sustained effort.",
        "5.2": "Students sometimes think BARR 'does the work' for them — it doesn't. It amplifies the effort you bring.",
        "5.3": "Students think alternates are just backups that don't matter — they are equally valid selections and should be chosen intentionally.",
        "5.4": "Students assume they should ask impressive-sounding questions for the panel. The best questions are honest ones.",
    }
    label = lesson["lesson_label"]
    return misconceptions.get(label, "Watch for students who conflate the vocabulary terms — clarify distinctions clearly.")


# ════════════════════════════════════════════════════════════════════════════
# MAIN BUILD
# ════════════════════════════════════════════════════════════════════════════

def link_btn(slide, text, x, y, w, h, fill_color, text_color, url, size=14):
    """
    Create a clickable button in one shape: fill + centered bold text + hyperlink.
    Using one shape (not rect+textbox) ensures clicks hit the hyperlink.
    """
    sh = rect(slide, x, y, w, h, fill_color)
    tf = sh.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    # Hyperlink on the text run (reliably clickable)
    run.hyperlink.address = url
    # Also wire the shape-level click action as a fallback
    _add_hyperlink_to_shape(sh, url)
    return sh


def _add_hyperlink_to_shape(shape, url):
    """
    Attach a click-through URL hyperlink to a shape using direct XML manipulation.
    python-pptx's click_action API works but only writes hlinkClick; this writes
    it directly into cNvPr so PowerPoint reliably opens the URL on click.
    """
    import lxml.etree as etree

    # Register the relationship on the slide part and get the rId
    rId = shape.part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )

    # Find the cNvPr element inside the shape's nvSpPr
    sp = shape._element
    # Presentationml namespace
    PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    cNvPr = sp.find(f".//{{{PML}}}cNvPr")
    if cNvPr is None:
        # Fallback: find via drawingml namespace (used in some shape types)
        DML2 = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
        cNvPr = sp.find(f".//{{{DML2}}}cNvPr")

    if cNvPr is not None:
        # Remove any existing hlinkClick child, then add fresh one
        for existing in cNvPr.findall(f"{{{DML}}}hlinkClick"):
            cNvPr.remove(existing)
        hlinkClick = etree.SubElement(cNvPr, f"{{{DML}}}hlinkClick")
        hlinkClick.set(f"{{{R}}}id", rId)
    else:
        # Last resort: use python-pptx's built-in click_action API
        shape.click_action.hyperlink.address = url


def video_slide(prs, unit, lesson):
    """Full dedicated slide for the YouTube video resource for this lesson."""
    yt = YOUTUBE.get(lesson["lesson_label"])
    if not yt:
        return
    color = uc(unit)

    s = new_slide(prs)
    bg(s, DARK)
    rect(s, 0, 0, 10, 0.09, color)
    rect(s, 0, 7.41, 10, 0.09, color)

    # Header chip
    chip = rect(s, 0.4, 0.22, 1.8, 0.34, color)
    tb(s, "▶  VIDEO", 0.42, 0.24, 1.76, 0.30, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Lesson label
    tb(s, f"Lesson {lesson['lesson_label']}", 2.4, 0.22, 7, 0.35, size=11, color=GOLD)

    # Main video title
    tb(s, yt["title"], 0.4, 0.72, 9.2, 1.3, size=26, bold=True, color=WHITE)

    # Channel + duration
    tb(s, f"{yt['channel']}  ·  approx. {yt['duration']}",
       0.4, 2.08, 9.2, 0.4, size=13, italic=True, color=GOLD)

    # Divider
    rect(s, 0.4, 2.6, 9.2, 0.03, color)

    # "Watch for" box
    rect(s, 0.4, 2.72, 9.2, 1.05, RGBColor(0x1A, 0x1A, 0x2E))
    tb(s, "AS YOU WATCH —", 0.6, 2.78, 2.2, 0.36, size=10, bold=True, color=color)
    tb(s, yt["watch_for"], 0.6, 3.1, 8.8, 0.62, size=13, color=WHITE)

    # Alt video suggestion — single linked shape
    rect(s, 0.4, 3.95, 9.2, 0.7, RGBColor(0x22, 0x22, 0x38))
    tb(s, "ALSO TRY →", 0.6, 4.0, 1.6, 0.28, size=9, bold=True, color=GOLD)
    alt_label = f'"{yt["alt_title"]}"  |  {yt["alt_channel"]}'
    link_btn(s, alt_label, 2.3, 4.0, 7.0, 0.56,
             RGBColor(0x22, 0x22, 0x38), GOLD, yt["alt_url"], size=11)

    # Big WATCH button — single shape with text run hyperlink
    link_btn(s, "▶  CLICK TO WATCH", 3.0, 4.9, 4.0, 0.82,
             color, WHITE, yt["url"], size=16)

    # Discussion prompt at bottom
    tb(s, "After watching: turn to a partner and answer the 'As You Watch' question in 60 seconds.",
       0.4, 5.9, 9.2, 0.4, size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

    notes(s,
          f"VIDEO RESOURCE — Lesson {lesson['lesson_label']}\n\n"
          f"Primary: \"{yt['title']}\" — {yt['channel']} ({yt['duration']})\n"
          f"URL: {yt['url']}\n\n"
          f"Alt option: \"{yt['alt_title']}\" — {yt['alt_channel']}\n"
          f"URL: {yt['alt_url']}\n\n"
          f"WATCH-FOR PROMPT: {yt['watch_for']}\n\n"
          f"FACILITATION:\n"
          f"• Click the button to open YouTube (internet required)\n"
          f"• Tell students the 'As You Watch' prompt BEFORE starting the video\n"
          f"• After the video: 60-second pair-share, then cold-call 1-2 students\n"
          f"• Note: if YouTube is blocked on school network, search the title in "
          f"  SchoolTube or Common Sense Media, or download ahead of time with "
          f"  a district-approved tool.\n\n"
          f"TIMING: This slide is optional — use when time allows after the reading "
          f"or as a hook before the content section.")


# ════════════════════════════════════════════════════════════════════════════
# MONDAY TEMPLATE — Look Back at the Week
# ════════════════════════════════════════════════════════════════════════════

MONDAY_COLOR = RGBColor(0x1A, 0x37, 0x6A)   # deep navy

def monday_slides(prs):
    """Three-slide Monday routine: agenda → tracker → notes review."""

    # ── Slide 1: Monday Check-In Title ──────────────────────────────────
    s = new_slide(prs)
    bg(s, MONDAY_COLOR)
    rect(s, 0, 0, 10, 0.09, GOLD)
    rect(s, 0, 7.41, 10, 0.09, GOLD)

    tb(s, "MONDAY CHECK-IN", 0.4, 0.5, 9.2, 0.5, size=14, bold=True,
       color=GOLD, align=PP_ALIGN.CENTER)
    tb(s, "Looking Back at Last Week", 0.4, 1.05, 9.2, 1.1, size=34, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.5, 2.3, 7, 0.04, GOLD)

    # 4-block agenda
    agenda = [
        ("📋", "Missing Work", "What's not turned in yet?"),
        ("📝", "Tests & Quizzes", "Review results + make-up plan"),
        ("📖", "Notes Review", "What did you learn in other classes last week?"),
        ("✅", "Weekly Win", "One thing you did well"),
    ]
    x_positions = [0.3, 2.75, 5.2, 7.65]
    for i, (icon, title, sub) in enumerate(agenda):
        x = x_positions[i]
        rect(s, x, 2.55, 2.2, 1.8, WHITE)
        rect(s, x, 2.55, 2.2, 0.5, GOLD)
        tb(s, icon, x + 0.05, 2.56, 0.5, 0.46, size=18, align=PP_ALIGN.CENTER)
        tb(s, title, x + 0.55, 2.6, 1.6, 0.4, size=11, bold=True, color=DARK)
        tb(s, sub, x + 0.1, 3.12, 2.0, 0.9, size=9, color=DARK, italic=True)

    tb(s, "Take out your Warrior Prep workbook — Monday Check-In page.",
       0.4, 4.55, 9.2, 0.45, size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s, "You have 3 minutes to complete the top section independently before we share out.",
       0.4, 5.05, 9.2, 0.4, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    timer_chip(s, 20, x=8.5, y=6.5)

    notes(s,
          "MONDAY CHECK-IN — AGENDA SLIDE\n\n"
          "OPENING ROUTINE (2-3 min):\n"
          "SAY: 'Good morning. Open your planner. Open Infinite Campus. "
          "Let's start where every good week starts — by looking back at last week.'\n\n"
          "PACE:\n"
          "• Missing Work Tracker: 5 min (students log, teacher circulates)\n"
          "• Test/Quiz Review: 8 min (debrief scores + make-up tracker)\n"
          "• Notes Review: 10 min (Cornell from other classes, pair-share)\n"
          "• Weekly Win: 5 min (whole class share, positive tone)\n"
          "• Exit + goal-set for the week: 2 min\n\n"
          "TEACHER TIP: Keep Monday Check-In to 30 minutes max. "
          "Use Infinite Campus on the projector to model checking grades together.")

    # ── Slide 2: Missing Work + Test Tracker ────────────────────────────
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Monday Check-In  ·  Missing Work & Test Review", MONDAY_COLOR)

    # Missing Work section
    rect(s, 0.2, 0.68, 4.6, 0.35, MONDAY_COLOR)
    tb(s, "📋  MISSING WORK TRACKER", 0.3, 0.71, 4.4, 0.28, size=11, bold=True, color=WHITE)

    col_headers = ["Class / Course", "Assignment", "Due Date", "Status"]
    col_widths   = [1.3, 1.5, 0.8, 0.85]
    y = 1.1
    xs = [0.25, 1.6, 3.15, 4.0]
    for i, (h, w) in enumerate(zip(col_headers, col_widths)):
        rect(s, xs[i], y, w, 0.28, LIGHT)
        tb(s, h, xs[i] + 0.04, y + 0.04, w - 0.08, 0.2, size=8, bold=True, color=DARK)
    y += 0.32
    for row in range(5):
        shade = CREAM if row % 2 == 0 else WHITE
        for i, (h, w) in enumerate(zip(col_headers, col_widths)):
            rect(s, xs[i], y, w, 0.34, shade)
        y += 0.36

    # Test/Quiz Review section
    rect(s, 5.1, 0.68, 4.65, 0.35, MONDAY_COLOR)
    tb(s, "📝  TEST & QUIZ REVIEW", 5.2, 0.71, 4.45, 0.28, size=11, bold=True, color=WHITE)

    t_headers = ["Class", "Test / Quiz", "Score", "Make-Up?"]
    t_widths   = [1.0, 1.6, 0.7, 1.2]
    t_xs       = [5.15, 6.2, 7.85, 8.6]
    y2 = 1.1
    for i, (h, w) in enumerate(zip(t_headers, t_widths)):
        rect(s, t_xs[i], y2, w, 0.28, LIGHT)
        tb(s, h, t_xs[i] + 0.04, y2 + 0.04, w - 0.08, 0.2, size=8, bold=True, color=DARK)
    y2 += 0.32
    for row in range(5):
        shade = CREAM if row % 2 == 0 else WHITE
        for i, (h, w) in enumerate(zip(t_headers, t_widths)):
            rect(s, t_xs[i], y2, w, 0.34, shade)
        y2 += 0.36

    # Make-up commitment box
    rect(s, 0.2, 3.7, 9.6, 0.85, RGBColor(0xFF, 0xF9, 0xC4))
    rect(s, 0.2, 3.7, 0.06, 0.85, GOLD)
    tb(s, "MAKE-UP PLAN:  I will complete _______________________  by  ___________  during  ___________________",
       0.35, 3.84, 9.3, 0.55, size=11, color=DARK)

    # Discussion prompts
    rect(s, 0.2, 4.65, 9.6, 0.32, MONDAY_COLOR)
    tb(s, "CLASS DISCUSSION", 0.35, 4.69, 3.0, 0.24, size=10, bold=True, color=GOLD)
    prompts = [
        "1. If you have missing work — what got in the way? What is your specific plan to fix it this week?",
        "2. How did last week's tests go? If a score surprised you (good or bad), what explains it?",
    ]
    y3 = 5.05
    for p in prompts:
        tb(s, p, 0.3, y3, 9.4, 0.44, size=10.5, color=DARK)
        y3 += 0.5

    notes(s,
          "MISSING WORK + TEST REVIEW SLIDE\n\n"
          "MISSING WORK (5 min):\n"
          "• Students open Infinite Campus and log any missing assignments in the tracker\n"
          "• Circulate — do NOT announce names publicly; work quietly one-on-one\n"
          "• If a student has 3+ missing, ask them to write a make-up plan in the box\n\n"
          "TEST REVIEW (8 min):\n"
          "• Students log any tests/quizzes from last week and their scores\n"
          "• If a score is below 70% in any class, ask: 'When is the retake? What will you do differently?'\n"
          "• Make-up tracker: help students schedule a specific time (lunch, advisory, after school)\n\n"
          "DISCUSSION (3-4 min):\n"
          "• Q1: Normalize missing work — shift to problem-solving, not shame\n"
          "• Q2: Frame test scores as data, not identity")

    # ── Slide 3: Notes Review + Weekly Win ──────────────────────────────
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Monday Check-In  ·  Notes Review & Weekly Win", MONDAY_COLOR)

    # Notes Review
    rect(s, 0.2, 0.68, 5.8, 0.35, MONDAY_COLOR)
    tb(s, "📖  NOTES REVIEW FROM OTHER CLASSES", 0.3, 0.71, 5.6, 0.28, size=11, bold=True, color=WHITE)

    directions = (
        "Cornell-style review:  Cover your notes.  Using just your cue column (left side), "
        "try to answer your own questions.  Then uncover and check.  "
        "Write a 1-sentence summary of the most important idea from each class."
    )
    tb(s, directions, 0.3, 1.1, 5.6, 0.75, size=10.5, color=DARK, italic=True)

    class_labels = ["Class 1:", "Class 2:", "Class 3:"]
    y4 = 1.95
    for cl in class_labels:
        rect(s, 0.25, y4, 5.65, 0.5, LIGHT)
        tb(s, cl, 0.35, y4 + 0.08, 0.75, 0.34, size=10, bold=True, color=DARK)
        tb(s, "Most important idea from last week: ________________________________",
           1.15, y4 + 0.1, 4.65, 0.3, size=9.5, italic=True, color=GRAY)
        y4 += 0.6

    # Pair-share prompt
    rect(s, 0.25, 3.78, 5.65, 0.68, RGBColor(0xE3, 0xF2, 0xFD))
    rect(s, 0.25, 3.78, 0.06, 0.68, MONDAY_COLOR)
    tb(s, "PAIR-SHARE: Tell your partner the most surprising thing you learned "
          "in any class last week. 60 seconds each.",
       0.38, 3.86, 5.4, 0.52, size=11, color=DARK, italic=True)

    # Weekly Win
    rect(s, 6.2, 0.68, 3.55, 0.35, GOLD)
    tb(s, "✅  WEEKLY WIN", 6.3, 0.71, 3.35, 0.28, size=11, bold=True, color=DARK)

    rect(s, 6.2, 1.1, 3.55, 3.38, CREAM)
    tb(s, "Write ONE thing you are proud of from last week.\n"
          "It can be academic, social, athletic, creative — anything real.",
       6.35, 1.18, 3.25, 0.7, size=10.5, color=DARK, italic=True)
    for line_y in [2.05, 2.5, 2.95, 3.35]:
        rect(s, 6.35, line_y, 3.2, 0.02, RGBColor(0xC0, 0xC0, 0xB0))

    tb(s, "SHARE OUT:", 6.3, 3.6, 1.5, 0.3, size=10, bold=True, color=GOLD)
    tb(s, "Who wants to share?", 7.85, 3.62, 1.8, 0.28, size=10, color=DARK, italic=True)

    # Goal for the week
    rect(s, 0.2, 4.6, 9.6, 0.35, MONDAY_COLOR)
    tb(s, "🎯  MY GOAL FOR THIS WEEK", 0.35, 4.63, 5.0, 0.28, size=11, bold=True, color=GOLD)
    rect(s, 0.2, 4.97, 9.6, 0.52, LIGHT)
    tb(s, "This week I will _______________________________________________________________________ by _____________",
       0.35, 5.06, 9.3, 0.36, size=11, color=DARK)

    rect(s, 0.2, 5.55, 9.6, 0.42, CREAM)
    tb(s, "Warrior Prep Word for the Week:  Stay on top of it before it gets on top of you.",
       0.35, 5.62, 9.3, 0.28, size=10.5, bold=True, italic=True, color=MONDAY_COLOR,
       align=PP_ALIGN.CENTER)

    notes(s,
          "NOTES REVIEW + WEEKLY WIN\n\n"
          "NOTES REVIEW (10 min):\n"
          "SAY: 'Open the Cornell notes from your hardest class last week. "
          "Cover your notes and try to answer your cue questions from memory. "
          "This is called retrieval practice — it triples how much you remember.'\n\n"
          "STRUCTURE:\n"
          "1. Students review notes silently (4 min)\n"
          "2. Pair-share: 'What was the most surprising thing you learned in any class?' (2 min)\n"
          "3. 2-3 students share with the class (2 min)\n\n"
          "WEEKLY WIN (5 min):\n"
          "• Ask 3-4 students to share — rotate so every student shares at least once per unit\n"
          "• Frame wins broadly: finishing a hard assignment, showing up on time, helping someone\n\n"
          "WEEK GOAL:\n"
          "• 1-2 minutes: students write a specific, achievable goal for the week\n"
          "• SAY: 'Put this in your planner so you see it every day this week.'")


# ════════════════════════════════════════════════════════════════════════════
# FRIDAY TEMPLATE — Look Ahead to Next Week
# ════════════════════════════════════════════════════════════════════════════

FRIDAY_COLOR = RGBColor(0x2E, 0x7D, 0x32)   # forest green

def friday_slides(prs):
    """Three-slide Friday routine: agenda → planner build → commitment."""

    # ── Slide 1: Friday Planning Title ──────────────────────────────────
    s = new_slide(prs)
    bg(s, FRIDAY_COLOR)
    rect(s, 0, 0, 10, 0.09, GOLD)
    rect(s, 0, 7.41, 10, 0.09, GOLD)

    tb(s, "FRIDAY PLANNING", 0.4, 0.5, 9.2, 0.5, size=14, bold=True,
       color=GOLD, align=PP_ALIGN.CENTER)
    tb(s, "Looking Ahead to Next Week", 0.4, 1.05, 9.2, 1.1, size=34, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.5, 2.3, 7, 0.04, GOLD)

    agenda = [
        ("📅", "Deadlines", "Exams, projects, big assignments"),
        ("🏃", "Activities", "Sports, clubs, personal plans"),
        ("🗓️", "Weekly Plan", "Build your day-by-day schedule"),
        ("💪", "Commitment", "One thing you will protect"),
    ]
    x_positions = [0.3, 2.75, 5.2, 7.65]
    for i, (icon, title, sub) in enumerate(agenda):
        x = x_positions[i]
        rect(s, x, 2.55, 2.2, 1.8, WHITE)
        rect(s, x, 2.55, 2.2, 0.5, GOLD)
        tb(s, icon, x + 0.05, 2.56, 0.5, 0.46, size=18, align=PP_ALIGN.CENTER)
        tb(s, title, x + 0.55, 2.6, 1.6, 0.4, size=11, bold=True, color=DARK)
        tb(s, sub, x + 0.1, 3.12, 2.0, 0.9, size=9, color=DARK, italic=True)

    tb(s, "Open your planner to next week. Open Infinite Campus. Let's build the week before it builds itself.",
       0.4, 4.55, 9.2, 0.45, size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    timer_chip(s, 20, x=8.5, y=6.5)

    notes(s,
          "FRIDAY PLANNING — AGENDA SLIDE\n\n"
          "OPENING (1-2 min):\n"
          "SAY: 'Good afternoon. We've made it to Friday. Before you pack up and disappear, "
          "let's do the one thing that separates students who have a good week from students who "
          "get surprised by one: we are going to plan.'\n\n"
          "PACE:\n"
          "• Upcoming Deadlines: 6 min (log all known tests, projects, due dates)\n"
          "• Activities Tracker: 5 min (personal + school, afternoon + evening)\n"
          "• Weekly Plan Builder: 8 min (block out homework time around activities)\n"
          "• Commitment Statement: 3 min (write + share)\n\n"
          "TEACHER TIP: Do this at a consistent time every Friday — students will "
          "start looking forward to it as a week-closure ritual.")

    # ── Slide 2: Deadlines + Activities ─────────────────────────────────
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Friday Planning  ·  Deadlines & Activities Next Week", FRIDAY_COLOR)

    # Upcoming Deadlines
    rect(s, 0.2, 0.68, 5.8, 0.35, FRIDAY_COLOR)
    tb(s, "📅  UPCOMING DEADLINES  (next 2 weeks)", 0.3, 0.71, 5.6, 0.28, size=11, bold=True, color=WHITE)

    d_headers = ["Class / Course", "Assignment / Test", "Due Date", "Points"]
    d_widths   = [1.25, 2.1, 1.0, 0.8]
    d_xs       = [0.25, 1.55, 3.7, 4.75]
    y = 1.1
    for i, (h, w) in enumerate(zip(d_headers, d_widths)):
        rect(s, d_xs[i], y, w, 0.28, LIGHT)
        tb(s, h, d_xs[i] + 0.04, y + 0.04, w - 0.08, 0.2, size=8, bold=True, color=DARK)
    y += 0.32
    for row in range(6):
        shade = CREAM if row % 2 == 0 else WHITE
        for i, w in enumerate(d_widths):
            rect(s, d_xs[i], y, w, 0.32, shade)
        y += 0.34

    # Priority flag
    rect(s, 0.25, 3.78, 5.65, 0.5, RGBColor(0xFF, 0xEB, 0xEE))
    rect(s, 0.25, 3.78, 0.06, 0.5, RED)
    tb(s, "HIGHEST PRIORITY this week: _______________________________________________",
       0.38, 3.9, 5.4, 0.28, size=10.5, color=DARK)

    # Activities column
    rect(s, 6.2, 0.68, 3.55, 0.35, FRIDAY_COLOR)
    tb(s, "🏃  ACTIVITIES NEXT WEEK", 6.3, 0.71, 3.35, 0.28, size=11, bold=True, color=WHITE)

    day_labels = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Weekend"]
    y_act = 1.1
    for day in day_labels:
        rect(s, 6.2, y_act, 3.55, 0.36, LIGHT if day_labels.index(day) % 2 == 0 else WHITE)
        tb(s, day + ":", 6.28, y_act + 0.06, 0.8, 0.24, size=9, bold=True, color=FRIDAY_COLOR)
        tb(s, "_________________________", 7.1, y_act + 0.07, 2.5, 0.22, size=9, color=GRAY)
        y_act += 0.38

    rect(s, 6.2, 3.42, 3.55, 0.85, RGBColor(0xE8, 0xF5, 0xE9))
    rect(s, 6.2, 3.42, 0.06, 0.85, FRIDAY_COLOR)
    tb(s, "Which activity needs the most preparation this week?\n"
          "Write it here: ________________________________",
       6.3, 3.5, 3.35, 0.68, size=9.5, color=DARK, italic=True)

    # Discussion
    rect(s, 0.2, 4.4, 9.6, 0.32, FRIDAY_COLOR)
    tb(s, "CLASS CHECK", 0.35, 4.44, 2.0, 0.24, size=10, bold=True, color=GOLD)
    for i, q in enumerate([
        "Does anyone have TWO or more major deadlines in the same week? Let's problem-solve.",
        "Does anyone have a day this week with ZERO time for homework? Let's build it in now."
    ]):
        tb(s, f"{i+1}. {q}", 0.3, 4.82 + i * 0.46, 9.4, 0.4, size=10.5, color=DARK)

    notes(s,
          "DEADLINES + ACTIVITIES SLIDE\n\n"
          "DEADLINES (6 min):\n"
          "• Students open Infinite Campus and check all upcoming assignments\n"
          "• They log anything due in the NEXT 2 WEEKS (not just next week)\n"
          "• They circle the single highest-priority item and write it in the red box\n"
          "• TEACHER: Scan the room. Look for students with 3+ things due same day\n\n"
          "ACTIVITIES (4 min):\n"
          "• Students log ALL scheduled activities — not just school ones\n"
          "• Include: sports practices, games, religious events, family events, work, babysitting\n"
          "• This is why students run out of time — invisible commitments\n\n"
          "CLASS CHECK (3 min):\n"
          "• Q1: Normalize deadline pile-ups. Help students prioritize and break tasks down\n"
          "• Q2: Find the 'dead days' where homework never happens — make a specific plan")

    # ── Slide 3: Weekly Plan Builder + Commitment ────────────────────────
    s = new_slide(prs)
    bg(s, WHITE)
    header_bar(s, "Friday Planning  ·  Build Your Week + Make Your Commitment", FRIDAY_COLOR)

    # Weekly plan grid
    rect(s, 0.2, 0.68, 6.9, 0.35, FRIDAY_COLOR)
    tb(s, "🗓️  MY WEEK AT A GLANCE — Build It Now", 0.3, 0.71, 6.7, 0.28,
       size=11, bold=True, color=WHITE)

    days = ["Mon", "Tue", "Wed", "Thu", "Fri"]
    time_blocks = ["After school\n(3–5 PM)", "Dinner break\n(5–6 PM)", "Evening\n(6–9 PM)", "Bedtime\n(9–10 PM)"]
    col_w = 1.2
    row_h = 0.56
    header_h = 0.3
    grid_x = 0.25
    grid_y = 1.1

    # Column headers (days)
    rect(s, grid_x, grid_y, 1.25, header_h, LIGHT)   # blank corner
    for j, day in enumerate(days):
        cx = grid_x + 1.3 + j * col_w
        rect(s, cx, grid_y, col_w - 0.04, header_h, FRIDAY_COLOR)
        tb(s, day, cx + 0.02, grid_y + 0.04, col_w - 0.08, header_h - 0.08,
           size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Row headers + cells
    for i, block in enumerate(time_blocks):
        ry = grid_y + header_h + i * row_h
        rect(s, grid_x, ry, 1.25, row_h - 0.02, LIGHT)
        tb(s, block, grid_x + 0.04, ry + 0.06, 1.16, row_h - 0.1, size=7.5, color=DARK)
        for j in range(len(days)):
            cx = grid_x + 1.3 + j * col_w
            shade = RGBColor(0xF1, 0xF8, 0xE9) if i % 2 == 0 else WHITE
            rect(s, cx, ry, col_w - 0.04, row_h - 0.02, shade)

    # Homework protection prompt
    rect(s, 0.25, grid_y + header_h + len(time_blocks) * row_h + 0.04, 7.0, 0.36, RGBColor(0xE8, 0xF5, 0xE9))
    tb(s, "⭐  Block your HOMEWORK TIME first, then fit everything else around it.",
       0.35, grid_y + header_h + len(time_blocks) * row_h + 0.1, 6.8, 0.24,
       size=10, bold=True, italic=True, color=FRIDAY_COLOR)

    # Commitment box
    rect(s, 7.35, 0.68, 2.4, 5.9, CREAM)
    rect(s, 7.35, 0.68, 2.4, 0.35, GOLD)
    tb(s, "💪 COMMITMENT", 7.45, 0.71, 2.2, 0.28, size=10, bold=True, color=DARK)

    commit_items = [
        ("I will protect:", "homework time from _____ to _____ every school day"),
        ("I will finish:", "_________________________________ by _________"),
        ("I will ask for help with:", "_________________________________ from _________"),
        ("I will sleep by:", "_______ PM on school nights"),
        ("I will NOT:", "let _________________ distract me from my work"),
    ]
    cy = 1.1
    for label, fill in commit_items:
        tb(s, label, 7.45, cy, 2.2, 0.24, size=8, bold=True, color=FRIDAY_COLOR)
        tb(s, fill, 7.45, cy + 0.26, 2.2, 0.44, size=8, italic=True, color=DARK)
        cy += 0.82

    # Closing
    rect(s, 0.2, 6.8, 9.6, 0.52, FRIDAY_COLOR)
    tb(s, "\"The best time to plan your week is before it starts.  You're doing that right now.  "
          "Have a great weekend, Warriors.\"",
       0.35, 6.88, 9.3, 0.36, size=10, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

    notes(s,
          "WEEKLY PLAN BUILDER + COMMITMENT\n\n"
          "PLAN BUILDER (8 min):\n"
          "SAY: 'Take your activities from the last slide and put them in the grid. "
          "Then block your homework time. Work backward from your deadlines.'\n\n"
          "FACILITATION:\n"
          "1. Students fill in activities first (2 min)\n"
          "2. Students block homework time — at least 90 min per school day (3 min)\n"
          "3. Students identify which evening is their hardest and make a plan (2 min)\n\n"
          "COMMITMENT (3 min):\n"
          "• Each student completes the commitment column\n"
          "• Ask 2-3 students to read ONE commitment aloud — makes it social\n"
          "• SAY: 'Take a photo of this page. Text it to someone who will hold you to it.'\n\n"
          "CLOSING:\n"
          "• Read the quote at the bottom aloud\n"
          "• SAY: 'Have a great weekend. We'll check in Monday to see how it went.'")


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    course_title_slide(prs)

    for unit in UNITS:
        unit_divider_slide(prs, unit)
        vocab_deep_dive_slide(prs, unit)
        pretest_slide(prs, unit)

        for lesson in unit["lessons"]:
            bell_ringer_slide(prs, unit, lesson)
            lesson_title_slide(prs, unit, lesson)
            reading_slides(prs, unit, lesson)
            video_slide(prs, unit, lesson)       # YouTube resource after reading
            story_slides(prs, unit, lesson)
            content_slides(prs, unit, lesson)
            exit_ticket_slide(prs, unit, lesson)

        posttest_slide(prs, unit)

    out = "output/warrior_prep_teacher_slides.pptx"
    prs.save(out)
    slide_count = len(prs.slides)
    print(f"PowerPoint saved → {out}  ({slide_count} slides)")

    # ── Weekly Routine Slide Deck ────────────────────────────────────────
    prs2 = Presentation()
    prs2.slide_width = Inches(10)
    prs2.slide_height = Inches(7.5)

    # Title
    s = new_slide(prs2)
    bg(s, DARK)
    rect(s, 0, 0, 10, 0.09, RED)
    rect(s, 0, 7.41, 10, 0.09, RED)
    tb(s, SCHOOL_NAME, 0.4, 1.2, 9.2, 0.55, size=18, color=GOLD, align=PP_ALIGN.CENTER)
    tb(s, "Weekly Routine Templates", 0.4, 1.85, 9.2, 1.2, size=38, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, "Monday Check-In  ·  Friday Planning", 0.4, 3.15, 9.2, 0.55,
       size=20, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    rect(s, 1.5, 3.85, 7, 0.04, GOLD)
    tb(s, "Use every Monday and Friday throughout the semester.",
       0.4, 4.0, 9.2, 0.4, size=14, color=WHITE, align=PP_ALIGN.CENTER)

    monday_slides(prs2)
    friday_slides(prs2)

    out2 = "output/warrior_prep_weekly_routine.pptx"
    prs2.save(out2)
    print(f"Weekly routine slides saved → {out2}  ({len(prs2.slides)} slides)")


if __name__ == "__main__":
    build()
